import os
from copy import deepcopy
from json import load, loads, dump
from multiprocessing import Process
from shutil import copy2, move, rmtree
from random import randint
from threading import Lock, Thread
from zipfile import ZipFile, ZIP_DEFLATED

import numpy as np
from moviepy import VideoFileClip
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pyaudio import PyAudio
from pydub import AudioSegment
from transformers import pipeline

import tkinter as tk
from tkinter import ttk, font
from tkinter import filedialog, messagebox, simpledialog
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from matplotlib.ticker import MultipleLocator


def speed_audio(audio: AudioSegment, factor: float) -> AudioSegment:
    """ 音频倍速
    :param audio: 音频对象
    :type audio: AudioSegment
    :param factor: 倍速
    :type factor: float
    :return: 开好倍速的音频
    :rtype: AudioSegment """
    # 原始帧率
    original_rate = audio.frame_rate
    # 调整帧率
    new_rate = int(original_rate * factor)
    audio_speed = audio._spawn(audio.raw_data, overrides={"frame_rate": new_rate})
    # 转换为原始帧率以保持音频长度正确
    audio_speed = audio_speed.set_frame_rate(original_rate)
    return audio_speed


class AudioNotLoadError(Exception):
    """ 未加载音频 """
    def __str__(self) -> str: return "Haven't loaded any audio file."

class AudioNotPredError(Exception):
    """ 未识别音频 """
    def __init__(self, name: str): self.name = name
    def __str__(self) -> str: return f"Haven't predicted the audio \"{self.name}\"."


class SpeechHandler:
    """
    音频处理
    ---
    语音转文字 & 文本对齐。此外还可以手动校准对齐信息。

    :param name: 音频文件名称
    :type name: str
    :param path: 音频路径(禁止修改)
    :type path: str
    :param info: 音频信息
    :type info: list[dict]
    :param audio: 加载的音频
    :type audio: AudioSegment
    :param already_load: 是否加载音频
    :type already_load: bool
    :param already_pred: 是否识别音频
    :type already_pred: bool

    :param __init__: 初始化(可能不加载模型)
    :param clean: 清理临时生成的音频文件
    :param cut: 切割对齐
    :param get_slice: 获取一段切片的音频
    :param get_text: 获取全部文本(以列表形式)
    :param get_stamp: 获取时间戳
    :param init_from: 从保存信息加载自己
    :param load: 加载音频
    :param load_config: 显式同步设置
    :param merge_above: 将指定索引切片合并到上一项
    :param predict: 识别文本与对齐时间戳
    :param reset: 重置信息
    :param save: 保存临时工作信息(含音频)
    :param video_load: 从视频中加载音频
    """

    def __init__(self, config: dict):
        """ 初始化模型
        :param config: 设置
        :type config: dict """
        self.reset()
        # self.workpath = os.getcwd()  # 工作路径 # 去 load_config 里看
        self.tmp_audio = []  # 临时音频文件存放位置
        # 加载设置
        self.load_config(config)
    
    def save(self, file_path: str):
        """ 打包项目内容
        :param file_path: 保存路径
        :type file_path: str """
        # 不使用 __getstate__ 和 __setstate__
        self._raise_null()
        # 临时路径
        work_path = os.path.join(self.tmp_path, f"{self.name}_pack{randint(1000, 10000)}")
        os.makedirs(work_path)
        # 打包信息记录
        pack_info = {
            "name": self.name,
            "info": self.info,
            "file": os.path.basename(self.origin), 
            "is_audio": self.is_audio, 
        }
        info_path = os.path.join(work_path, "data.json")
        with open(info_path, mode="w", encoding="utf-8") as inff:
            dump(pack_info, inff)
        # 进行打包
        zip_path = os.path.join(work_path, "package.zip")
        with ZipFile(zip_path, "w", ZIP_DEFLATED) as zipf:
            # 打包信息存储
            zipf.write(info_path, "data.json")
            # 音/视频打包
            zipf.write(self.origin, pack_info["file"])
        # 移动包
        move(zip_path, file_path)
        # 清理临时文件
        rmtree(work_path)
    
    def init_from(self, file_path: str):
        """ 从项目打包中恢复
        :param file_path: 打包路径
        :type file_path: str """
        with ZipFile(file_path, "r", ZIP_DEFLATED) as zipf:
            # 获取信息
            with zipf.open("data.json", mode="r") as inff:
                info = loads(inff.read().decode("utf-8"))
            # 将原音/视频拉入临时文件夹
            spq_path = os.path.join(self.tmp_path, info["file"])
            try:
                # zipf.extract(info["file"], spq_path)  # <- 这是错的
                zipf.extract(info["file"], self.tmp_path)
            except:
                os.makedirs(self.tmp_path)
                zipf.extract(info["file"], self.tmp_path)
            self.tmp_audio.append(spq_path)
        # 载入信息
        # self.reset()
        if info["is_audio"]:
            self.load(spq_path, bl_pred=False)
        else:
            self.video_load(spq_path, bl_pred=False)
        self.info = info["info"]  # 载入对齐信息
        self.already_pred = True  # 已经有了对齐
    
    def load_config(self, config: dict):
        """ 加载设置
        :param config: 设置
        :type config: dict """
        self.workpath = config["workpath"]  # 工作路径
        self.tmp_path = os.path.join(self.workpath, config["temp_path"])  # 临时文件路径
        self.model_name = config["model"]  # 模型位置
        self.load_all_time = bool(config["load_all_time"])  # 启动就加载模型

        if self.load_all_time:
            self.unload_model()  # 保险一点
            self.load_model()
    
    def load_model(self):
        """ 加载模型 """
        # 加载模型
        # 完全没看出直接用成品有任何技术含量...
        self.pipe = pipeline(
            "automatic-speech-recognition",
            model=self.model_name,
            chunk_length_s=30,  # 启用分块处理
            stride_length_s=6,
            device="cpu",
        )
    
    def unload_model(self):
        """ 取消加载模型 """
        try:
            del self.pipe
        finally:
            self.pipe = None
    
    def predict(self):
        """ 预测 """
        self._raise_null(bl_pred=False)
        self.already_pred = True
        if not self.load_all_time:
            self.load_model()
        self.info = self.pipe(self.path, batch_size=8, return_timestamps=True)["chunks"]
        if not self.load_all_time:
            self.unload_model()
        for sen in self.info:
            sen["text"] = sen["text"].strip()

    def reset(self):
        """ 重置音频信息 """
        self.already_load = self.already_pred = False
        self.name = ""
        self.path = ""
        self.origin = ""  # 初始文件地址(不论音频还是视频)(打包时使用)
        self.is_audio: bool = None  # 初始文件类型
        self.info = []  # [{"timestamp": (), "text": str}, ]
        self.audio = None  # 音频文件
    
    def load(self, audio_path: str, bl_pred: bool=True):
        """ 载入音频信息
        :param audio_path: 音频路径
        :type audio_path: str
        :param bl_pred: 是否识别
        :type bl_pred: bool """
        self.reset()
        self.origin = audio_path
        self.is_audio = True
        self.path = audio_path
        self.name = os.path.basename(self.path)
        self.audio = AudioSegment.from_file(self.path)
        # 格式转换
        if not os.path.basename(self.name) in [".mp3", ".wav", ".flac"]:
            path = os.path.join(self.tmp_path, self.name + ".wav")
            try:
                self.audio.export(path, format="wav")
            except:
                os.makedirs(self.tmp_path)
                self.audio.export(path, format="wav")
            finally:
                self.path = path
                self.tmp_audio.append(self.path)
        self.already_load = True
        # 识别音频并对齐文本
        if bl_pred:
            self.predict()
            # ↑↑↑ 战场遗址: 2025/3 未经处理的 torch-directml 只会告诉我 !*n 消耗 >3h

    def video_load(self, video_path: str, bl_pred: bool=True):
        """ 载入音频信息
        :param video_path: 视频路径
        :type video_path: str
        :param bl_pred: 是否识别
        :type bl_pred: bool """
        self.reset()
        self.origin = video_path
        self.is_audio = False
        # 临时音频文件位置
        self.name = f"{os.path.basename(video_path)}"
        self.path = os.path.join(self.tmp_path, self.name + ".wav")
        self.tmp_audio.append(self.path)
        # 保存
        with VideoFileClip(video_path) as video:
            try:
                video.audio.write_audiofile(self.path)
            except:
                os.makedirs(self.tmp_path)
                video.audio.write_audiofile(self.path)
        # 加载音频
        self.audio = AudioSegment.from_file(self.path)  # 不是 .wav 就要安装 ffmpeg
        self.already_load = True
        # 识别音频并对齐文本
        if bl_pred:
            self.predict()

    def cut(self, key: int, stamp: float, text_index: int):
        """ 切割一段对齐信息
        :param key: 句子索引
        :type key: int
        :param stamp: 切割点时间戳
        :type stamp: float
        :param text_index: 文本切割位
        :type text_index: int """
        self._raise_null()
        sen1 = {"text": self.info[key]["text"][:text_index], "timestamp": (self.info[key]["timestamp"][0], stamp)}
        sen2 = {"text": self.info[key]["text"][text_index:], "timestamp": (stamp, self.info[key]["timestamp"][1])}
        self.info[key] = sen1
        self.info.insert(key + 1, sen2)

    def merge_above(self, key: int):
        """ 将一段切片合并到上一项
        :param key: 该段切片的索引
        :type key: int """
        self._raise_null()
        # 追加进上一项
        self.info[key - 1]["text"] += self.info[key]["text"]
        self.info[key - 1]["timestamp"] = (self.info[key - 1]["timestamp"][0], 
                                           self.info[key]["timestamp"][1])
        del self.info[key]
    
    def get_slice(self, key_from: int, key_to: int=None):
        """ 获取一段音频切片
        :param key_from: 该段切片的索引
        :type key_from: int
        :param key_to: 该段切片的索引
        :type key_to: int = None (默认只有 key_from 这一段切片)
        :return: 音频切片
        :rtype: AudioSegment """
        self._raise_null()
        if key_to == None:
            key_to = key_from
        return self.audio[self.info[key_from]["timestamp"][0]*1000 : self.info[key_to]["timestamp"][1]*1000]

    def get_text(self) -> list:
        """ 获取全部文本
        :return: 全部文本以及文本戳
        :rtype: [text: str, stamp: list[int]] """
        self._raise_null()
        text = ""
        stamp = []
        for sen in self.info:
            stamp += len(text), 
            text += " " + sen["text"]
        return text[1:], stamp
        # return [sen["text"] for sen in self.info]

    def get_stamp(self) -> list[float]:
        """ 获取所有分段时间戳
        :return: 时间戳(不含开头 0.00s)
        :rtype: list[float] """
        self._raise_null()
        return [sen["timestamp"][1] for sen in self.info][:-1]

    def clean(self):
        """ 清理临时文件 """
        for files in self.tmp_audio:
            try:  # 无论如何，窗口必须顺利停止，所以遇到问题直接忽略
                os.remove(files)
            except Exception as err:
                print(err)

    def _raise_null(self, bl_pred: bool=True):
        """ 未加载报错 """
        if not self.already_load:
            raise AudioNotLoadError()
        if bl_pred and not self.already_pred:
            raise AudioNotPredError(self.name)

    def generate(self, slides: list[dict], ppt_path: str):
        """ PPT 生成
        :param slides: 每张 PPT 的信息 (下文介绍)
        :type slides: list[  
                dict['title': str, 'type': 'single', 'key': int, 'text': bool, 'single': bool, 'repeat': (time: int>=0, sep: float), 'speed': ('single': float=None, 'repeat': float=None) or factor: float or None]
                    # 单句 #
                    # text: 是否追加句子显示及其动画, repeat.time: 重复次数(time=0, 不添加重复的控件), repeat.sep: 间隔时间 #
                or dict['title': str, 'type': 'long', 'text': bool, 'side': [from, end: int], 'speed': float or None]
                    # 长段 #
                    # side: 起点和终点(其一为 None 表示到尽头) #
            ]
        :param ppt_path: 导出的 PPT 的位置
        :type ppt_path: str """
        self._raise_null()
        # 工作路径
        work_path = os.path.join(self.tmp_path, f"{self.name}_ppt{randint(1000, 10000)}")
        try:
            rmtree(work_path)
        except:
            pass
        try:
            os.makedirs(work_path, exist_ok=True)  # 创建目录
        except:
            os.makedirs(self.tmp_path)
            os.makedirs(work_path, exist_ok=True)  # 创建目录
        finally:
            os.makedirs(os.path.join(work_path, "audio"), exist_ok=True)  # 创建 audio 目录
        # 拷贝必要文件
        copy2(os.path.join(self.workpath, "icon", "audio.png"), work_path)  # 音频图标
        icon = "audio.png"
        # 转到工作路径
        os.chdir(work_path)

        # PPT 设计参数
        wd, ht = Pt(960), Pt(540)  # 幻灯片尺寸
        lf_x, cen_x, mid_x, all_w = Pt(36), Pt(500), Pt(450), Pt(888)  # 左侧及中间统一线 & 横向跨越
        sep_y, row1_y = Pt(100), Pt(124)  # 标题-内容分界线 & 小标题行
        tt2_w, tt2_h, aud_y, aud_wh = Pt(240), Pt(36), Pt(160), Pt(72)  # 音频按钮
        ans_h = Pt(300)  # 答案呈现部分正文
        sep_c, ft_c = RGBColor(128, 128, 128), RGBColor(64, 64, 64)  # 分界线 & 字体颜色
        tt2_s, con_s = Pt(22), Pt(26)  # 小标题字号，正文字号
        tt_n, con_n, ln_p = "微软雅黑", "Times New Roman", 1.2  # 字体和行距
        # 创建 PPT
        prs = Presentation()
        prs.slide_width = wd  # 尺寸
        prs.slide_height = ht
        slide_layout = prs.slide_layouts[1]  # 正文布局
        # 逐页添加
        cnt = -1
        for config in slides:
            cnt += 1
            # 识别类型
            single: bool = config["type"] == "single"  # 单张为 True
            # 如果为长文本，先识别起点和终点
            if not single:
                if config["side"][0] == None:
                    config["side"][0] = 0
                if config["side"][1] == None:
                    config["side"][1] = len(self.info)
                    
            # 添加新页面
            slide = prs.slides.add_slide(slide_layout)

            # 标题
            if not config["title"]:
                if single:  # 格式: 句子 X
                    config["title"] = f"句子 {config['key'] + 1}"
                else:  # 格式: 长文本 X~X
                    config["title"] = f"长句 {config['side'][0] + 1}~{config['side'][1] + 1}"
            # 添加标题
            title = slide.shapes.title
            title.text = config["title"]
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # 左对齐
            title.text_frame.paragraphs[0].font.name = tt_n

            # 标题-内容 分界线
            line = slide.shapes.add_shape(autoshape_type_id=1,  # 直线
                left=lf_x, top=sep_y, width=all_w, height=Pt(0))
            line.line.color.rgb = sep_c

            # mode=Single -> 单句录音部分
            if single:
                # 获取原始切片
                audio_sen = self.get_slice(config["key"])  # 本句切片
                # 设置倍速信息
                if config["speed"] == None:
                    speed = {"speed": False}  # 不开倍速
                elif type(config["speed"]) == float:
                    speed = {"speed": True, "both": False, "factor": config["speed"]}  # 一起倍速
                    audio_sen = speed_audio(audio_sen, speed["factor"])  # 覆写音频
                else:  # 不同倍速
                    speed = {"speed": True, "both": True, "single": False if config["speed"][0] else ({"audio": speed_audio(audio_sen, config["speed"][0]), "factor": config["speed"][0]}),\
                            "repeat": False if config["speed"][1] else {"audio": speed_audio(audio_sen, config["speed"][1]), "factor": config["speed"][1]}}

                # 提取音频保存
                if config["repeat"][0]:
                    # 重复片段
                    audio_repeat = os.path.join("audio", name := f"page{cnt}sen{config["key"]}_repeat.wav")
                    silence = AudioSegment.silent(duration=config["repeat"][1] * 1000)  # 静音片段
                    audio = AudioSegment.empty()  # 最终音频
                    # 选择音频
                    if speed["speed"] and speed["both"]:
                        audio_base = speed["repeat"]["audio"] if speed["repeat"] else audio_sen
                    else:
                        audio_base = audio_sen
                    # 合成音频
                    for i in range(config["repeat"][0]):
                        audio += audio_base
                        if i < config["repeat"][0]:
                            audio += silence
                    audio.export(audio_repeat, format="wav")
                if config["single"]:
                    # 单个片段
                    audio_single = os.path.join("audio", name := f"page{cnt}sen{config["key"]}_single.wav")
                    if speed["speed"] and speed["both"]:
                        (speed["single"]["audio"] if speed["single"] else audio_sen).export(os.path.join(work_path, audio_single))
                    else:
                        audio_sen.export(audio_single, format="wav")
                
                # 决定排版模式
                if config["repeat"][0] and config["single"]:
                    # 双栏排版
                    # left -> 小标题
                    left_content = slide.shapes.add_textbox(left=lf_x, top=row1_y, width=tt2_w, height=tt2_h)
                    tt2_l = left_content.text_frame.add_paragraph()
                    tt2_l.text = "单遍录音" if ((not speed["speed"]) or (speed["speed"] and speed["both"] and not speed["single"])) else f"单遍录音 (x{(speed['factor'] if not speed['both'] else round(speed['single']['factor'], 1))})"
                    tt2_l.font.size = tt2_s
                    tt2_l.font.color.rgb = ft_c
                    tt2_l.font.name = tt_n
                    # left -> 音频
                    slide.shapes.add_movie(movie_file=audio_single, poster_frame_image=icon, mime_type="audio/wav", 
                        left=lf_x, top=aud_y, width=aud_wh, height=aud_wh)
                    # right -> 小标题
                    right_content = slide.shapes.add_textbox(left=cen_x, top=row1_y, width=tt2_w, height=tt2_h)
                    tt2_r = right_content.text_frame.add_paragraph()
                    tt2_r.text = f"{config['repeat'][0]} 遍录音" if ((not speed["speed"]) or (speed["speed"] and speed["both"] and not speed["repeat"])) else f"{config['repeat'][1]} 遍录音 (x{(speed['factor'] if not speed['both'] else round(speed['repeat']['factor']), 1)})"
                    tt2_r.font.size = tt2_s
                    tt2_r.font.color.rgb = ft_c
                    tt2_r.font.name = tt_n
                    # right -> 音频
                    slide.shapes.add_movie(movie_file=audio_repeat, poster_frame_image=icon, mime_type="audio/wav", 
                        left=cen_x, top=aud_y, width=aud_wh, height=aud_wh)
                else:
                    # 单栏排版
                    content = slide.shapes.add_textbox(left=mid_x, top=row1_y, width=tt2_w, height=tt2_h)
                    tt2 = content.text_frame.add_paragraph()
                    tt2.text = ("单遍录音" if not speed["speed"] else f"单遍录音 (x{round(speed['factor'], 1)})") if config["single"] else (f"{config['repeat'][0]} 遍录音" if not speed["speed"] else f"{config['repeat'][0]} 遍录音 (x{round(speed['factor'], 1)})")
                    tt2.font.size = tt2_s
                    tt2.font.color.rgb = ft_c
                    tt2.font.name = tt_n
                    # tt2.alignment = PP_ALIGN.CENTER
                    # 音频
                    slide.shapes.add_movie(movie_file=audio_single if config["single"] else audio_repeat, poster_frame_image=icon, mime_type="audio/wav", 
                        left=mid_x, top=aud_y, width=aud_wh, height=aud_wh)
            
            # mode=Long -> 长文本正文部分
            else:
                # 获取原始切片
                audio_sen = self.get_slice(*config["side"])  # 多句切片切片
                # 设置倍速信息
                if config["speed"] != None:
                    audio_sen = speed_audio(audio_sen, config["speed"])
                # 保存音频
                audio_sen.export(audio_path := os.path.join("audio", name := f"page{cnt}.wav"), format="wav")
                # 标题排版
                content = slide.shapes.add_textbox(left=mid_x, top=row1_y, width=tt2_w, height=tt2_h)
                tt2 = content.text_frame.add_paragraph()
                tt2.text = "原速度录音" if config["speed"] == None else f"x{round(config['speed'], 1)} 倍速"
                tt2.font.size = tt2_s
                tt2.font.color.rgb = ft_c
                tt2.font.name = tt_n
                # tt2.alignment = PP_ALIGN.CENTER
                # 图标
                slide.shapes.add_movie(poster_frame_image=icon, movie_file=audio_path, mime_type="audio/wav", 
                    left=mid_x, top=aud_y, width=aud_wh, height=aud_wh)

            # 答案呈现部分
            if config["text"]:
                cnt += 1
                # 添加新页面
                slide = prs.slides.add_slide(slide_layout)
                # 添加标题
                title = slide.shapes.title
                title.text = config["title"]
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # 左对齐
                title.text_frame.paragraphs[0].font.name = tt_n
                # 标题-内容 分界线
                line = slide.shapes.add_shape(autoshape_type_id=1,  # 直线
                    left=lf_x, top=sep_y, width=all_w, height=Pt(0))
                line.line.color.rgb = sep_c

                # 添加答案文本
                textbox = slide.shapes.add_textbox(left=lf_x, top=row1_y, width=all_w, height=ans_h)
                text_frame = textbox.text_frame
                text_frame.word_wrap = True  # 自动换行
                p = text_frame.add_paragraph()
                if single:
                    p.text = self.info[config["key"]]["text"].strip()
                else:
                    p.text = "".join([f"{sen['text'].strip()} " for sen in self.info[config["side"][0] : config["side"][1]]])
                p.font.size = con_s  # 大字号
                p.font.name = con_n
                p.line_spacing = ln_p

        # 保存 PPT 项目
        prs.save(ppt_path)  # 保存
        # 压缩整个文件夹并打包
        # with ZipFile(zip_path := os.path.join(export_path, f"{name} 压缩包.zip"), "w", ZIP_DEFLATED) as zipf:
        #     # 音频打包
        #     for file in os.listdir(audio_path := os.path.join(work_path, "audio")):
        #         file_path = os.path.join(audio_path, file)
        #         zipf.write(file_path, os.path.join("audio", file))
        #     # ppt 打包
        #     zipf.write(ppt_path, name)
        # 切回原路径
        os.chdir(self.workpath)
        # 清理垃圾
        rmtree(work_path)


class Window:
    """
    本工具图形化界面(tkinter 版)
    ---

    :param __init__: 初始化窗口(不显示)，预载管理器和模型
    :param mainloop: 显示窗口并进入窗口主循环
    :param close: 关闭窗口并自动释放内存
    """

    class ScrollableFrame(tk.Frame):
            """
            滚动容器
            ---
            基于 tk.Frame 构建。  
            初版代码由 DeepSeek-R1 生成，经本人 Debug 并调整。

            注意，加入控件时请使用 ScrollabelFrame.container ，而非 ScrollabelFrame !
            """

            def __init__(self, parent, scroll_direction="both", **kwargs):
                """ 创建一个通用的滚动容器
                :param scroll_direction: 滚动方向 ("horizontal", "vertical", "both")
                :type scroll_diretion: str """
                tk.Frame.__init__(self, parent, **kwargs)
                
                # 初始化配置
                match scroll_direction:
                    case "both":
                        self.scroll_direction = {"x": True, "y": True}
                    case "vertical":
                        self.scroll_direction = {"x": False, "y": True}
                    case "horizontal":
                        self.scroll_direction = {"x": True, "y": False}
                self._setup_widgets()
                self._bind_events()

            def _setup_widgets(self):
                """ 创建界面元素 """
                # 创建Canvas和滚动条
                self.canvas = tk.Canvas(self, highlightthickness=0)
                self.v_scroll = ttk.Scrollbar(self, orient="vertical")
                self.h_scroll = ttk.Scrollbar(self, orient="horizontal")

                # 配置滚动条
                self.canvas.configure(
                    yscrollcommand=self.v_scroll.set if self.scroll_direction["y"] else None,
                    xscrollcommand=self.h_scroll.set if self.scroll_direction["x"] else None
                )
                # ↑↑↑ 战场遗址: 2025/4/3 很喜欢互换的 x 和 y 啊。

                # 创建内部容器
                self.container = tk.Frame(self.canvas)
                self.canvas.create_window((0, 0), window=self.container, anchor="nw", tags="frame_container")

                # 布局元素
                self._arrange_layout()

            def _arrange_layout(self):
                """ 根据滚动方向布局组件 """
                if self.scroll_direction["y"]:
                    self.v_scroll.pack(side="right", fill="y")
                if self.scroll_direction["x"]:
                    self.h_scroll.pack(side="bottom", fill="x")
                
                self.canvas.pack(side="left", fill="both", expand=True)

                # 配置滚动条命令
                if self.scroll_direction["y"]:
                    self.v_scroll.config(command=self.canvas.yview)
                if self.scroll_direction["x"]:
                    self.h_scroll.config(command=self.canvas.xview)

            def _bind_events(self):
                """ 绑定必要的事件 """
                # 鼠标滚动
                self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)
                self.canvas.bind_all("<Shift-MouseWheel>", self._on_shift_mousewheel)
                # 容器尺寸变化事件
                self.canvas.bind("<Configure>", self._on_canvas_configure)
                self.container.bind("<Configure>", self._on_container_configure)

            def _on_mousewheel(self, event):
                """处理垂直滚动"""
                if self.scroll_direction["y"]:
                    self.canvas.yview_scroll(-1*(event.delta//120), "units")

            def _on_shift_mousewheel(self, event):
                """处理水平滚动( Shift+滚轮)"""
                if self.scroll_direction["x"]:
                    self.canvas.xview_scroll(-1*(event.delta//120), "units")

            def _on_container_configure(self, event):
                """ 当内部容器尺寸变化时更新滚动区域 """
                self.canvas.configure(scrollregion=self.canvas.bbox("all"))
            
            def _on_canvas_configure(self, event):
                """ 当画布尺寸变化时调整内部容器尺寸 """
                self.canvas.itemconfigure("frame_container", width=event.width)


    class Menu(tk.Menu):
        """ 顶部菜单 """

        def __init__(self, monitor, *args, **kwargs):
            """ 初始化菜单
            :param monitor: 程序主体类 -> 沟通信息及相互调用
            :type monitor: Window """
            tk.Menu.__init__(self, monitor.root, *args, **kwargs)
            self.monitor = monitor
            self.monitor.root.config(menu=self)

            # 菜单构建
            self.add_command(label="新窗口(N)", underline=4, command=self.monitor.new_window)
            
            self.openMenu = tk.Menu(self.monitor.root, tearoff=False)
            self.add_cascade(label="打开(O)", underline=3, menu=self.openMenu)
            self.openMenu.add_command(label="打开音频(A)", underline=5, command=self.monitor.open_audio)
            self.openMenu.add_command(label="打开视频(V)", underline=5, command=self.monitor.open_video)
            self.openMenu.add_separator()
            self.openMenu.add_command(label="还原存档(R)", underline=5, command=self.monitor.load_work)

            self.exportMenu = tk.Menu(self.monitor.root, tearoff=False)
            self.add_cascade(label="导出(E)", underline=3, menu=self.exportMenu, state="disabled")
            self.exportMenu.add_command(label="生成听写PPT(P)", underline=8, command=self.monitor.generate)
            self.exportMenu.add_command(label="导出音频切片(A)", underline=7, command=self.monitor.export_slice)
            self.exportMenu.add_command(label="导出文本.txt(T)", underline=9, command=self.monitor.export_txt)
            self.exportMenu.add_separator()
            self.exportMenu.add_command(label="保存存档(A)", underline=5, command=self.monitor.save_work)
            
            self.add_command(label="设置(S)", underline=3, command=self.monitor.config_board)
            self.add_command(label="帮助(H)", underline=3, command=self.monitor.help_board)


    class TopBar(tk.Frame):
        """ 顶部信息和快速预览条 """

        def __init__(self, monitor, *args, **kwargs):
            """ 初始化顶部状态条
            :param monitor: 程序主体类 -> 沟通信息及相互调用
            :type monitor: Window """
            tk.Frame.__init__(self, monitor.root, *args, **kwargs)
            self.monitor = monitor

            self.fileLb1 = tk.Label(self, text="文件:")
            self.fileLb2 = tk.Label(self, font=self.monitor.fonts["text3_title"])
            self.set_info()

            self.textBt = ttk.Button(self, text="全文", width=5, command=self._load_viewer, state="disabled")
            self.generateBt = ttk.Button(self, text="生成 PPT", width=9, state="disabled", command=self.monitor.generate)
            self.saveBt = ttk.Button(self, text="保存", state="disabled", width=5, command=self.monitor.save_work)

            # 排列控件
            self.fileLb1.pack(side="left", padx=1)
            self.fileLb2.pack(side="left")
            self.saveBt.pack(side="right", padx=2)
            self.generateBt.pack(side="right")
            self.textBt.pack(side="right", padx=2)
        
        def set_info(self):
            """ 加载信息 """
            self.fileLb2.config(text=self.monitor.audio.name if self.monitor.audio.name else "无")

        def activate(self):
            """ 激活按钮 """
            self.textBt.config(state="normal")
            self.generateBt.config(state="normal")
            self.saveBt.config(state="normal")

        def _load_viewer(self):
            """ 加载全文阅览器 """
            viewer = self.monitor.TextViewer(self.monitor.audBrd)


    class SentenceBoard(tk.Frame):
        """ 句子显示板
        ---
        
        :param clear: 清空所有句子的显示(可从指定位置开始)
        :param load: 加载句子(可从指定位置开始)
        :param highlight: 手动高亮
        """

        class SentenceDetail(tk.Frame):
            """
            句子细节及操作展示
            ---

            :param __init__: 初始化控件和信息
            :param highlight: 高亮
            :param unhighlight: 取消高亮
            :param set_info: 根据索引设置信息
            :param copy: 复制本句文本
            :param export: 导出本句音频
            """

            def __init__(self, monitor, parent, key: int, **kwargs):
                """ 初始化控件信息
                :param monitor: 程序主体类 -> 沟通信息及相互调用
                :type monitor: Window
                :param parent: 父控件
                :param key: 本条句子的索引
                :type key: int """
                tk.Frame.__init__(self, parent, **kwargs)
                self.monitor = monitor
                self.parent = parent

                # 静态信息展示控件
                self.keyLb = tk.Label(self, font=self.monitor.fonts["text1_hold"])
                self.stampLb = tk.Label(self, fg=self.monitor.colors["gray2"])
                self.senLb = tk.Label(self, font=self.monitor.fonts["text2_para"], fg=self.monitor.colors["gray1"], anchor="nw", justify="left")
                # ↑↑↑ 战场遗址: 2025/3 我们把 self 写成了 parent，获得了 >1h DEBUG
                # 获取并刷新信息
                self.set_info(key)

                # 右击菜单
                self.menu = tk.Menu(self, tearoff=False)
                self.menu.add_command(label="跳转(J)", underline=3, command=lambda: self.monitor.audBrd.control.goto_sen(self.key))
                self.menu.add_command(label="复制(C)", underline=3, command=self.copy)
                self.menu.add_command(label="导出音频(E)", underline=5, command=self._ask_export)
                self.menu.add_separator()
                self.menu.add_command(label="向上合并(M)", underline=5, command=self.merge_above)
                self.menu.add_command(label="编辑(E)", underline=3, command=self._edit)

                # 高亮显示控件
                self.hlFrm = tk.Frame(self)
                self.jumpBt = ttk.Button(self.hlFrm, text="跳转", width=5, command=lambda: self.monitor.audBrd.control.goto_sen(self.key))
                self.sep1 = ttk.Separator(self.hlFrm, orient="vertical")
                self.mergeBt = ttk.Button(self.hlFrm, text="向上合并", width=9, command=self.merge_above)
                if self.key == 0:
                    self.mergeBt.config(state="disabled")
                self.editBt = ttk.Button(self.hlFrm, text="编辑", width=5, command=self._edit)
                self.sep2 = ttk.Separator(self.hlFrm, orient="vertical")
                self.exportBt = ttk.Button(self.hlFrm, text="导出音频", width=9, command=self._ask_export)
                self.copyBt = ttk.Button(self.hlFrm, text="复制文本", width=9, command=self.copy)

                # 排布静态控件
                self.keyLb.pack(side="top", anchor="nw")
                self.stampLb.pack(side="top", anchor="nw")
                self.senLb.pack(side="top", anchor="nw", fill="x", expand=True)
                # 排布高亮控件
                self.jumpBt.pack(side="left", padx=(6, 0))
                self.sep1.pack(side="left", padx=8, fill="y")
                self.mergeBt.pack(side="left", padx=(0, 4))
                self.editBt.pack(side="left")
                self.sep2.pack(side="left", padx=8, fill="y")
                self.exportBt.pack(side="left")
                self.copyBt.pack(side="left", padx=(4, 0))

                # 高亮绑定
                self.config(highlightbackground=self.monitor.colors["gray3"], highlightthickness=1)
                for widget in self.winfo_children() + [self, ]:
                    widget.bind("<Button-1>", func=self.highlight)  # 高亮绑定
                    widget.bind("<Button-3>", self._post_menu)  # 菜单绑定
                # self.bind_all("<Button-1>", func=self.highlight)
                # self.bind_all("<Button-3>", self._post_menu)
                # ↑↑↑ 战场遗址: 2025/4/12 不可使用 bind_all，因为这会 bind 窗口里其他的控件
                self.senLb.bind("<Configure>", self._update_wrap)  # 自动换行绑定
                # map(lambda widget: widget.bind("<Button-1>", func=self.highlight), [\
                #     self, self.keyLb, self.stampLb, self.senLb, ])

            def _edit(self):
                """ 唤出编辑窗口 """
                edit = self.monitor.senBrd.EditToplevel(self.monitor.senBrd)
            
            def set_info(self, key: int):
                """ 设置信息
                :param key: 本条句子的索引
                :type key: int """
                # 记录信息
                self.key = key
                self.info = self.monitor.audio.info[key]
                # 显示信息
                self.keyLb.config(text=f"句子 {self.key + 1}")
                self.stampLb.config(text=f"{self.info["timestamp"][0]:.2f}s  →  {self.info["timestamp"][1]:.2f}s")
                self.senLb.config(text=self.info["text"])
            
            def unhighlight(self):
                """ 取消高亮 """
                self.hlFrm.pack_forget()
                self.config(highlightbackground=self.monitor.colors["gray3"], highlightthickness=1)
            
            def highlight(self, event):
                """ 被高亮 """
                self.hlFrm.pack(side="top", anchor="nw", fill="x")
                self.config(highlightbackground=self.monitor.colors["firebrick1"], highlightthickness=2)
                # 选中高亮传递信息
                self.monitor.senBrd.update_highlight(self.key)  # 执行选中
            
            def _update_wrap(self, event):
                """ 自动换行 """
                self.senLb.config(wraplength=event.width)

            def merge_above(self):
                """ 向上合并 """
                self.monitor.audio.merge_above(self.key)
                self.monitor.senBrd.load(self.key - 1)
                self.monitor.audBrd.control.update_stamp()

            def _post_menu(self, event):
                """ 唤出菜单 """
                self.menu.post(event.x_root, event.y_root)

            def copy(self):
                """ 复制本句文本 """
                self.clipboard_clear()
                self.clipboard_append(self.info["text"])
                self.update()

            def export(self, path: str):
                """ 导出音频文件
                :param path: 导出文件位置
                :type path: str """
                audio = self.monitor.audio.get_slice(self.key)
                type = os.path.splitext(path)
                audio.export(path, format=type[1:])

            def _ask_export(self):
                """ 询问位置并保存 """
                path = filedialog.asksaveasfilename(title=f"保存句子 {self.key + 1} 切片", defaultextension=".wav", initialfile=f"句子-{self.key + 1}.wav", filetypes=(("音频文件", "*.mp3 *.wav *.m4a *.flac"), ))
                self.export(path)


        class SelectMonitor:
            """
            选中项管理器
            ---

            :param select: 当前选中项
            :type select: int
            :param last: 上一个选中项
            :type last: int
            
            :param index: 获取当前序列号 + 1 的字符串
            :param change: 设置新的选中序列号，同时返回上一个选中项
            """

            def __init__(self, key: int=0):
                """ 初始化信息
                :param key: 初始索引
                :type key: int = 0 """
                self.select = key  # 当前选中索引
                self.last = key  # 上一个序列
            
            def index(self) -> str:
                """ 获取序号
                :return: 从 1 开始计数的序号
                :rtype: str """
                return str(self.select + 1)

            def change(self, key: int) -> int:
                """ 更换选中
                :param key: 新的索引
                :type key: int
                :return: 上一个索引
                :rtype: int """
                self.last, self.select = self.select, key
                return self.last


        class EditToplevel:
            """ 编辑 - 独立窗口 """
            
            def __init__(self, parent):
                """ 初始化一个切割窗口
                :param parent: 父类的父类
                :type parent: SentenceBoard """
                self.parent = parent
                self._load()

                # 窗口构建
                self.root = tk.Toplevel(self.parent)
                self.root.title(f"句子 {self.key} 编辑")
                self.root.resizable(False, False)

                # 控件及变量设置
                self.fromDb = tk.DoubleVar(self.root, self.side[0])
                self.fromLb = tk.Label(self.root, text=f"起始:")
                self.fromEt = ttk.Entry(self.root, textvariable=self.fromDb)
                self.endDb = tk.DoubleVar(self.root, self.side[1])
                self.endLb = tk.Label(self.root, text=f"终点:")
                self.endEt = ttk.Entry(self.root, textvariable=self.endDb)
                self.sep1 = ttk.Separator(self.root)
                self.txtTt = tk.Text(self.root, width=40, height=10)
                self.txtTt.insert("1.0", self.text)
                self.txtScr = ttk.Scrollbar(self.root, orient="vertical", command=self.txtTt.yview)
                self.txtTt.config(yscrollcommand=self.txtScr.set)
                self.sep2 = ttk.Separator(self.root)
                self.okBt = ttk.Button(self.root, text="确定", width=5, command=self.set)
                self.cancelBt = ttk.Button(self.root, text="取消", width=5, command=self.root.destroy)

                # 排布控件
                self.fromLb.grid(row=0, column=0, sticky="w")
                self.fromEt.grid(row=0, column=1, sticky="w")
                self.endLb.grid(row=1, column=0, sticky="w")
                self.endEt.grid(row=1, column=1, sticky="w")
                self.sep1.grid(row=2, column=0, columnspan=4, sticky="we")
                self.txtTt.grid(row=3, column=0, columnspan=3, sticky="nswe")
                self.txtScr.grid(row=3, column=3, sticky="ns")
                self.sep2.grid(row=4, column=0, columnspan=4, sticky="we")
                self.okBt.grid(row=5, column=1, sticky="e")
                self.cancelBt.grid(row=5, column=2, sticky="e")

                # 执行窗口
                self.root.mainloop()

            def _load(self):
                """ 载入信息 """
                self.key = self.parent.select.select
                self.side = self.parent.monitor.audio.info[self.key]["timestamp"]
                self.text = self.parent.monitor.audio.info[self.key]["text"]
            
            def set(self):
                """ 设置 """
                self.side = (self.fromDb.get(), self.endDb.get())
                self.text = self.txtTt.get("0.0", "end")
                if (bl_k1 := self.key == 0) or self.parent.monitor.audio.info[self.key - 1]["timestamp"][0] < self.side[0]:
                    if (bl_k2 := self.key == len(self.parent.monitor.audio.info) - 1) or self.side[1] < self.parent.monitor.audBrd.control.duration:
                        self.parent.monitor.audio.info[self.key] = {"text": self.text, "timestamp": self.side}
                        if not bl_k1:
                            self.parent.monitor.audio.info[self.key - 1]["timestamp"] = (self.parent.monitor.audio.info[self.key - 1]["timestamp"][0], self.side[1])
                        if not bl_k2:
                            self.parent.monitor.audio.info[self.key + 1]["timastamp"] = (self.side[0], self.parent.monitor.audio.info[self.key + 1]["timestamp"][1])
                        # 刷新 SentenceBoard 显示
                        self.parent.load(key_from=self.key - 1 if self.key > 0 else 0)  # 亲测后面的代码可以运行
                        self.parent.highlight(self.key)
                        # 刷新波形图
                        self.parent.monitor.audBrd.control.update_stamp()
                        # 关掉
                        self.root.destroy()
                        return
                # 不要乱输好吗
                messagebox.showwarning(title="警告", message="时间戳设置有误 !\n请检查与上下文时间戳是否有冲突。")


        def __init__(self, monitor, *args, **kwargs):
            """ 初始化句子板
            :param monitor: 程序主体类 -> 沟通信息及相互调用
            :type monitor: Window """
            tk.Frame.__init__(self, monitor.root, *args, **kwargs)
            self.monitor = monitor

            self.sentence = []  # 所有句子的容器
            self.senFrm = self.monitor.ScrollableFrame(self, scroll_direction="vertical")

            self.select = self.SelectMonitor()  # 选中管理器

            # 排列控件
            self.senFrm.pack(fill="both", padx=(4, 0), expand=True)
        
        def clear(self, key_from: int=0):
            """ 清除句子
            :param key_from: 起始句子
            :type key_from: int = 0 """
            # 清除选中
            if key_from >= self.select.select:
                self.select.change(0)
            # 删除控件
            for sen in self.sentence[key_from:]:
                sen.destroy()
            self.sentence = self.sentence[:key_from]

        
        def load(self, key_from: int=0):
            """ 加载句子
            :param key_from: 起始句子
            :type key_from: int = 0 """
            if self.sentence:
                self.clear(key_from=key_from)
            # 逐个添加
            for key in range(key_from, len(self.monitor.audio.info)):
                self.sentence += self.SentenceDetail(self.monitor, self.senFrm.container, key), 
                self.sentence[-1].pack(side="top", anchor="nw", fill="x", padx=2, pady=3)

        def update_highlight(self, key: int):
            """ 刷新选中(程序调用)
            :param key: 选中的索引
            :type key: int """
            # 同步选择信息并取消上一个高亮
            if (last := self.select.change(key)) != key:
                self.sentence[last].unhighlight()

        def highlight(self, key: int):
            """ 手动高亮
            :param key: 高亮的句子索引
            :type key: int """
            self.sentence[key].highlight(None)


    class AudioBoard(tk.Frame):
        """ 音频操作板 """

        class Preview(tk.Frame):
            """ 波形图预览
            ---
            
            :param load: 加载音频和时间戳
            :param reload_stamp: 重新加载时间戳
            :param goto: 跳转至指定时间
            """
           
            def __init__(self, parent, *args, **kwargs):
                """ 初始化波形图预览界面 """
                tk.Frame.__init__(self, parent, *args, **kwargs)

                # 基本信息设置
                self.width: int  # 界面宽度：秒
                self.place = 0.00  # 当前播放位置

                # 预览界面控件
                self.fig, self.ax = plt.subplots(figsize=(8, 3), facecolor="none")
                self.fig.subplots_adjust(left=0.02, right=0.98, top=0.85, bottom=0.1)  # 最大化
                self.canvas = FigureCanvasTkAgg(self.fig, master=self)  # 画布
                self.canvas.get_tk_widget().pack(fill="both", expand=True)

                # 样式配置
                self.ax.set_facecolor('none')
                self.ax.yaxis.set_visible(False)  # 隐藏纵轴
                self.ax.spines['top'].set_visible(False)
                self.ax.spines['right'].set_visible(False)
                self.ax.spines['bottom'].set_visible(False)
                self.ax.spines['left'].set_visible(False)
                self.ax.xaxis.set_ticks_position('top')  # 横轴刻度显示在顶部

                # 绑定尺寸变化
                self.bind("<Configure>", self._on_configure)
            
            def load(self, audio: AudioSegment, stamp: list[float]):
                """ 加载音频
                :param audio: 音频
                :type audio: AudioSegment
                :param stamp: 时间戳
                :type stamp: list[float] """
                self.duration = len(audio) / 1000  # 音频时长
                self.stamp = stamp  # 时间戳

                # 波形图预处理
                samples = np.array(audio.get_array_of_samples())  # 采样
                sample_rate = audio.frame_rate  # 采样率
                sample_count = len(samples)  # 样本数
                time_axis = np.linspace(0, self.duration, sample_count)  # 时间刻度设置
                # 完成采样(降低采样率)
                ds_factor = sample_rate // 100
                self.samples = samples[::ds_factor]  # 采样
                self.times = time_axis[::ds_factor]  # 对应时间刻度

                # 渲染波形图
                self._render()
            
            def reload_stamp(self, stamp: list[float]):
                """ 重新加载时间戳
                :param stamp: 时间戳
                :type stamp: list[float] """
                self.stamp = stamp
                self._render()

            def _render(self):
                """ 重新渲染波形图 """
                self.ax.clear()
                # 绘制波形图竖线
                markerline, stemlines, baseline = self.ax.stem(self.times, self.samples, linefmt="#404040", markerfmt=' ', basefmt=' ')
                plt.setp(stemlines, 'linewidth', 0.5)
                # 绘制时间戳
                for mark in self.stamp:
                    self.ax.axvline(mark, color="#FFD700", alpha=0.7, linewidth=1.2)
                # 绘制播放指针
                self.now = self.ax.axvline(0, color="red", linewidth=1.6)
                # 刻度设置
                self.ax.xaxis.set_major_locator(MultipleLocator(1.0))  # 主刻度每1秒
                self.ax.xaxis.set_minor_locator(MultipleLocator(0.1))  # 次刻度每0.1秒
                self.ax.tick_params(axis="x", which="both", direction="out", labelsize=12, length=5, width=1)

                # 跳转至原位置
                self.goto()

            def goto(self, place: float=None):
                """ 跳转至时间
                :param place: 目标位置
                :type place: float = self.place """
                if place:
                    self.place = place
                # 起始位置计算
                left = self.place - (half := self.width / 2)  # 左侧位置
                right = self.place + half  # 右侧位置
                # 绘制波形图
                self.ax.set_xlim(left, right)
                self.ax.figure.canvas.draw_idle()
                # 绘制播放指针
                self.now.set_xdata([self.place] * 2)
            
            def _on_configure(self, event):
                """ 适应尺寸变化 """
                self.width = event.width // 256

        
        class Control(tk.Frame):
            """ 播放控制部分
            ---

            :param place: 当前时间
            :type place: float
            :param is_playing: 是否在播放
            :type is_playing: bool
            :param duration: 音频总时长
            :type duration: float

            :param load: 加载音频和时间戳
            :param updata_stamp: 更新时间戳
            :param sentence: 查查当前位于哪个句子
            :param ms_place: 返回 分:秒 的时间
            :param widget_state: 调整按钮禁用和启用状态
            :param play: 播放
            :param pause: 暂停
            :param repeat: 重复
            :param goto: 跳转至指定位置
            :param goto_sen: 跳转至指定句子
            """
            
            def __init__(self, parent, configs: dict, *args, **kwargs):
                """ 初始化播放控制器
                :param parent: 父类
                :type parent: AudioBoard
                :param configs: 播放设置
                :type configs: dict """
                tk.Frame.__init__(self, parent, *args, **kwargs)
                self.parent = parent  # 我们操纵父类而不是最顶上的 Window
                self.configs = configs

                self.lock = Lock()
                # 控件定义
                self._setup_widget()
                self.widget_state(False)  # 初始可用状态: False
                
            def load(self, audio: AudioSegment, stamp: list[float]):
                """ 加载音频
                :param audio: 音频
                :type audio: AudioSegment
                :param stamp: 时间戳
                :type stamp: list[float] """
                # 音频信息
                self.audio = audio
                self.stamp = stamp
                # 播放信息 <- 为什么不扔进 __init__ ???
                self.place = 0.00  # 播放位置
                self.chunk_size = 1024
                self.pause_flag = True  # 是否暂停
                self.is_playing = False  # 是否暂停
                self.is_jumping = False  # 正在跳转
                self.goto_data_ptr = 0
                self.goto_place = 0
                self.play_thread = None
                self.stream = None

                # 波形图预载
                self.parent.preview.load(self.audio, self.stamp)
                # 音频就绪
                self._ready_audio()

            def _ready_audio(self):
                """ 准备音频，应对后面的播放 """
                # 初始化 PyAudio
                self.p = PyAudio()
                # 获取音频参数
                self.duration = len(self.audio) / 1000
                self.format = self.p.get_format_from_width(self.audio.sample_width)
                # self.channels = self.audio.channels
                self.rate = self.audio.frame_rate
                # 将音频转换为原始数据
                self.raw_data = self.audio.raw_data
                self.data_ptr = 0  # 当前播放数据指针

            def _play_thread(self):
                """ 播放线程 """
                while self.data_ptr < len(self.raw_data) and self.is_playing:
                    if self.pause_flag:
                        break
                    
                    end_ptr = self.data_ptr + self.chunk_size
                    chunk = self.raw_data[self.data_ptr : end_ptr]
                    self.stream.write(chunk)
                    self.data_ptr += self.chunk_size
                    
                    # 更新播放位置（毫秒）
                    self.place = (self.data_ptr / (self.audio.frame_rate 
                                * self.audio.sample_width 
                                * self.audio.channels)) 
                    self.parent.preview.goto(self.place)  # 更新波形图显示

                    # 跳转
                    if self.is_jumping:
                        # 同步数据
                        self.data_ptr = self.goto_data_ptr
                        self.place = self.goto_place
                        self.is_jumping = False
                        continue
                    
                self.stream.stop_stream()

                # 播完就暂停
                if not self.pause_flag:
                    self.data_ptr = len(self.raw_data)
                    self.place = self.duration
                    self.parent.preview.goto(self.place)  # 播放位置更新
                    self.grid_playBt()
                # 必须停
                self.pause_flag = False
                self.is_playing = False

            def play(self, bl_bt: bool=True):
                """ 从当前位置开始播放
                :param bl_bt: 是否更改暂停/播放控件
                :type bl_bt: bool """
                # 如果流存在且已停止，需要重新创建
                if self.stream and self.stream.is_stopped():
                    self.stream.close()
                    self.stream = None

                if self.stream == None:
                    # 创建音频流
                    self.stream = self.p.open(format=self.format, channels=self.audio.channels, rate=self.rate, output=True)
                # 转换为字节位置
                byte_pos = int(self.place * self.audio.frame_rate * self.audio.sample_width)
                # byte_pos = int(self.place * 1000 * self.audio.frame_rate * self.audio.sample_width)
                # print({"byte_pos": byte_pos, "data_ptr": self.data_ptr})
                # ↑↑↑ 战场遗址: 2025/4/4 15:44 结束 >2h DEBUG
                self.data_ptr = max(0, min(byte_pos, len(self.raw_data)))
                # 播完就重来
                if self.data_ptr >= len(self.raw_data):
                    self.data_ptr = 0
                    self.place = 0
                
                # 创建播放线程
                self.is_playing = True
                self.pause_flag = False
                self.play_thread = Thread(target=self._play_thread, daemon=True)  # 线程守护，关闭窗口即关闭进程，同时关掉这个线程。
                self.play_thread.start()
                # 刷新播放按钮
                if bl_bt:
                    self.grid_pauseBt()

            def pause(self, bl_bt: bool=True):
                """ 暂停播放
                :param bl_bt: 是否更改暂停/播放控件
                :type bl_bt: bool """
                # 停止音频流
                self.lock.acquire()
                if self.stream and self.stream.is_active():
                    self.pause_flag = True
                    self.lock.release()
                # 刷新播放按钮
                if bl_bt:
                    self.grid_playBt()

            def _setup_widget(self):
                """ 控件定义 """
                # 加载图片
                self.images = ["skip-backward", "fast-backward", "play", "pause", "fast-forward", "skip-forward", "repeat"]
                self.images = [Image.open(f"./icon/{img}.png").resize([45] * 2) for img in self.images]
                self.images = [ImageTk.PhotoImage(img) for img in self.images]
                # ↑↑↑ 战场遗址: 2025/4/4 有了抗体标记的 bug，即使过了两年也能秒 'v'
                # 控制按钮
                self.bkSenBt = ttk.Button(self, image=self.images[0], command=self._sen_backward)
                self.bkBt = ttk.Button(self,    image=self.images[1], command=self._backward)
                self.playBt = ttk.Button(self,  image=self.images[2], command=self.play)
                self.pauseBt = ttk.Button(self, image=self.images[3], command=self.pause)
                self.fwBt = ttk.Button(self,    image=self.images[4], command=self._forward)
                self.fwSenBt = ttk.Button(self, image=self.images[5], command=self._sen_forward)
                self.repeatBt = ttk.Button(self,image=self.images[6], command=self.repeat)
                # 排布控制按钮
                cnt = 0
                for widget in self.bkSenBt, self.bkBt, self.playBt, self.fwBt, self.fwSenBt, self.repeatBt:
                    widget.grid(row=0, column=cnt)
                    cnt += 1
                self.grid_columnconfigure([i for i in range(6)], weight=1)
            
            def _sen_backward(self):
                """ 回退一句 """
                self.goto_sen(self.sentence() - 1)

            def _sen_forward(self):
                """ 快进一句 """
                self.goto_sen(self.sentence() + 1)
            
            def _backward(self):
                """ 快退 """
                self.goto(self.place - self.configs["backward"])

            def _forward(self):
                """ 快进 """
                self.goto(self.place + self.configs["forward"])

            def goto(self, place: float):
                """ 跳转到指定位置播放
                :param place: 指定位置
                :type place: float """
                self.lock.acquire()  # 上！锁！
                # 跳完怎么办？那就跳完！
                if place >= self.duration:
                    self.goto_place = self.duration
                # 跳转。喜欢对子线程动手脚的小朋友们有福了
                self.goto_place = max(0.00, place)
                self.goto_data_ptr = int(place * self.audio.frame_rate * self.audio.sample_width * self.audio.channels)
                self.is_jumping = True
                # 我开锁，这个线程也许还没跳转。但我知道，它最终只能走向跳转。
                self.lock.release()

            def goto_sen(self, sen: int):
                """ 前往指定句子
                :param sen: 指定句子
                :type sen: int """
                if sen <= 0:
                    self.repeat()
                elif sen >= len(self.stamp):
                    self.goto(self.duration)
                else:
                    self.goto(self.stamp[sen - 1])
    
            def repeat(self):
                """ 重新播放 """
                self.goto(0.00)

            def grid_playBt(self):
                """ 显示播放按钮 """
                self.pauseBt.grid_forget()
                self.playBt.grid(row=0, column=2)
            
            def grid_pauseBt(self):
                """ 显示暂停按钮 """
                self.playBt.grid_forget()
                self.pauseBt.grid(row=0, column=2)
            
            def widget_state(self, able: bool=None):
                """ 激活或禁用按钮
                :param able: 是(True)否(False)激活
                :type able: bool = None (翻转当前状态) """
                # 状态设置
                if able == None:
                    state = not self.state
                else:
                    state = "normal" if able else "disabled"
                # 状态应用
                for widget in self.bkSenBt, self.bkBt, self.playBt, self.pauseBt, self.fwBt, self.fwSenBt, self.repeatBt:
                    widget.config(state=state)
                self.state = state

            def update_stamp(self):
                """ 更新时间戳 """
                self.stamp = self.parent.monitor.audio.get_stamp()
                self.parent.preview.reload_stamp(self.stamp)
            
            def ms_place(self) -> list[int]:
                """ 返回分秒为单位的时间
                :return: 分 : 秒
                :rtype: [int, int] """
                min = int(self.place) // 60 
                sec = int(self.place) % 60 
                return min, sec

            def sentence(self) -> int:
                """ 返回当前正在播放的句子
                :return: 当前正在播放的句子
                :rtype: int """
                for sen_cnt in range(len(self.stamp)):
                    if self.stamp[sen_cnt] >= self.place:
                        return sen_cnt
                return len(self.stamp)


        class CutToplevel:
            """ 切割 - 独立窗口 """

            def __init__(self, parent):
                """ 初始化一个切割窗口
                :param parent: 父类
                :type parent: AudioBoard """
                self.parent = parent
                self._pause_for_info()

                # 窗口构建
                self.root = tk.Toplevel(self.parent)
                self.root.title(f"句子 {self.key} 切割")
                self.root.resizable(False, False)

                # 控件及变量设置
                self.fromLb = tk.Label(self.root, text=f"起始:  {self.side[0]}s")
                self.endLb = tk.Label(self.root, text=f"终点:  {self.side[1]}s")
                self.cutLb = tk.Label(self.root, text="切割点:")
                self.cutDb = tk.DoubleVar(self.root, value=self.stamp)
                self.cutEt = ttk.Entry(self.root, textvariable=self.cutDb)
                self.sep1 = ttk.Separator(self.root)
                self.txtLb = tk.Label(self.root, text="请在下方文本框内找出切割处，输入 <cut> 以标注。\n除此之外请不要做出其他任何更改，否则会影响裁切！", justify="left", anchor="nw")
                self.txtTt = tk.Text(self.root, width=40, height=10)
                self.txtTt.insert("1.0", self.text)
                self.txtScr = ttk.Scrollbar(self.root, orient="vertical", command=self.txtTt.yview)
                self.txtTt.config(yscrollcommand=self.txtScr.set)
                self.sep2 = ttk.Separator(self.root)
                self.okBt = ttk.Button(self.root, text="确定", width=5, command=self.set)
                self.cancelBt = ttk.Button(self.root, text="取消", width=5, command=self.root.destroy)

                # 排布控件
                self.fromLb.grid(row=0, column=0, columnspan=4, sticky="w")
                self.endLb.grid(row=1, column=0, columnspan=4, sticky="w")
                self.cutLb.grid(row=2, column=0, sticky="w")
                self.cutEt.grid(row=2, column=1, sticky="w")
                self.sep1.grid(row=3, column=0, columnspan=4, sticky="we")
                self.txtLb.grid(row=4, column=0, columnspan=4, sticky="nw")
                self.txtTt.grid(row=5, column=0, columnspan=3, sticky="nswe")
                self.txtScr.grid(row=5, column=3, sticky="ns")
                self.sep2.grid(row=6, column=0, columnspan=4, sticky="we")
                self.okBt.grid(row=7, column=1, sticky="e")
                self.cancelBt.grid(row=7, column=2, sticky="e")

                # 执行窗口
                self.root.mainloop()
            
            def set(self):
                """ 设置并应用 """
                # 获取设置
                self.text = self.txtTt.get("0.0", "end")
                self.stamp = self.cutDb.get()
                # 应用设置
                try:
                    text_index = self.text.index("<cut>")
                    if self.text[text_index - 1] == " ":
                        text_index -= 1
                except:
                    messagebox.showwarning(title="警告", message="未能找到字符 \"<cut>\" !\n请重新输入")
                else:
                    if self.side[0] < self.stamp < self.side[1]:
                        # 同步信息
                        self.parent.monitor.audio.cut(self.key, self.stamp, text_index)
                        # 刷新波形图 & 播放器
                        self.parent.control.update_stamp()
                        # 刷新 SentenceBoard 显示
                        self.parent.monitor.senBrd.load(key_from=self.key)
                        # 关掉窗口
                        self.root.destroy()
                        return
                    else:
                        messagebox.showwarning(title="警告", message="切割位点范围不在当前句子中 !\n请重新设置")
            
            def _pause_for_info(self):
                """ 暂停播放并载入信息 """
                if self.parent.control.is_playing:
                    self.parent.control.pause()
                self.key = self.parent.control.sentence()
                self.stamp = self.parent.control.place  # 当前位点
                self.side = self.parent.monitor.audio.info[self.key]["timestamp"]
                self.text = self.parent.monitor.audio.info[self.key]["text"]


        def __init__(self, monitor, *args, **kwargs):
            """ 初始化音频操作板
            :param monitor: 程序主体类 -> 沟通信息及相互调用
            :type monitor: Window """
            tk.Frame.__init__(self, monitor.root, *args, **kwargs)
            self.monitor = monitor

            # 区域划分
            self.preview = self.Preview(self)
            self.control = self.Control(self, configs=self.monitor.configs)
            self.prompt = tk.Frame(self)

            # 其他信息和操作
            self.sep = ttk.Separator(self, orient="vertical")
            self.jumpBt = ttk.Button(self.prompt, text="选中", width=5, command=self._select_now, state="disabled")
            self.cutBt = ttk.Button(self.prompt, text="切割", width=5, command=self._ask_for_cut, state="disabled")
            self.placeLb = tk.Label(self.prompt, text="位置:  _:_")
            self.senLb = tk.Label(self.prompt, text="句子:  _")

            # 排布控件
            self.preview.grid(row=0, column=0, sticky="nswe", columnspan=3)
            self.control.grid(row=1, column=0, sticky="we")
            self.sep.grid(row=1, column=1, sticky="ns")
            self.prompt.grid(row=1, column=2, stick="we")
            self.jumpBt.grid(row=0, column=0, sticky="nw")
            self.cutBt.grid(row=1, column=0, sticky="nw")
            self.placeLb.grid(row=0, column=1, sticky="nw", padx=(16, 0))
            self.senLb.grid(row=1, column=1, sticky="nw", padx=(16, 0))

            # 刷新绑定
            self.after(self.monitor.configs["update"], lambda: self.update_label())
        
        def load(self):
            """ 加载音频 """
            self.control.load(self.monitor.audio.audio, self.monitor.audio.get_stamp())
        
        def update_label(self):
            """ 按钮刷新 """
            try:
                self.placeLb.config(text=f"位置:  {(ms_place := self.control.ms_place())[0]} : {ms_place[1] if ms_place[1] >= 10 else ' ' + str(ms_place[1])}")
                self.senLb.config(text=f"句子:  {self.control.sentence() + 1}")
            except:
                pass
            self.after(self.monitor.configs["update"], lambda: self.update_label())
        
        def _ask_for_cut(self):
            """ 唤出切割窗口 """
            cutToplevel = self.CutToplevel(self)

        def _select_now(self):
            """ SentenceBoard 选中当前句子 """
            self.monitor.senBrd.highlight(self.control.sentence())

        def activate(self):
            """ 激活控件 """
            self.jumpBt.config(state="normal")
            self.cutBt.config(state="normal")
            self.control.widget_state(True)


    class TextViewer:
        """
        文本同步阅览窗口
        ---
        与 AudioBoard 一同播放音频，本窗口保持高亮当前句子

        :param load: 重载文本
        :param update: 手动刷新定位
        """

        def __init__(self, parent):
            """ 初始化文本阅览窗口
            :param parent: 父控件
            :type parent: AudioBoard """
            self.parent = parent
            self.text = ""
            self.stamp = []
            self.stop_load = False  # 停止加载标签
            # 窗口设置
            self.root = tk.Toplevel(self.parent)
            self.root.title("文本同步阅览")
            self.root.minsize(500, 400)
            # 控件设置
            self._setup_widget()
            self._layout_widget()
            self.load()
            # 运行窗口
            self.root.mainloop()
        
        def load(self):
            """ 加载文本 """
            # 删除标签
            for cnt in range(len(self.stamp)):
                self.txtBrd.tag_delete(f"sen_{cnt}")
            # 载入信息
            self.last = 0
            self.first = True  # 初次载入
            self.text_update = True  # 文本刷新
            self.text, self.stamp = self.parent.monitor.audio.get_text()  # 文本和文本戳
            stamp = self.stamp + [len(self.text)]
            # 文本输入
            self.txtBrd.config(state="normal")
            self.txtBrd.delete("1.0", "end")
            for cnt in range(len(stamp) - 1):
                self.txtBrd.insert("end", self.text[stamp[cnt] : stamp[cnt + 1]])
                self.txtBrd.tag_add(f"sen_{cnt}", f"1.{stamp[cnt]}", f"1.{stamp[cnt + 1]}")  # 标签设置
            self.txtBrd.config(state="disabled")
            # 刷新位置和显示
            # self.update()
            self.root.after(self.parent.monitor.configs["update"], lambda: self.update())
        
        def update(self):
            """ 刷新位置和显示 """
            if self.stop_load:
                self.stop_load = False
                return
            try:
                # 当前句子位置
                self.key = self.parent.control.sentence()
                # 刷新控件
                self._update_widget()
            except:
                pass
            finally:
                self.root.after(self.parent.monitor.configs["update"], lambda: self.update())
        
        def _setup_widget(self):
            """ 设置控件 """
            self.txtFrm = tk.Frame(self.root)
            self.txtBrd = tk.Text(self.txtFrm, font=self.parent.monitor.fonts["text4_context"], state="disabled", width=24, height=12)
            self.scrll = ttk.Scrollbar(self.txtFrm, command=self.txtBrd.yview)
            self.txtBrd.config(yscrollcommand=self.scrll.set)
            self.sep1 = ttk.Separator(self.root)
            self.dataFrm = tk.Frame(self.root)
            self.wordLb = tk.Label(self.dataFrm)  # 字数统计
            self.sep2 = ttk.Separator(self.dataFrm, orient="vertical")
            self.placeLb = tk.Label(self.dataFrm)  # 当前位置 & 句数统计
            self.sep3 = ttk.Separator(self.dataFrm, orient="vertical")
            self.loadBt = ttk.Button(self.dataFrm, text="刷新", width=5, command=self._reload)
        
        def _layout_widget(self):
            """ 展示控件 """
            self.txtFrm.pack(side="top", fill="both", expand=True)
            self.sep1.pack(side="top", fill="x")
            self.dataFrm.pack(side="bottom", fill="x")
            self.txtBrd.pack(side="left", fill="both", expand=True)
            self.scrll.pack(side="right", fill="y")
            self.placeLb.grid(row=0, column=0, sticky="w")
            self.sep2.grid(row=0, column=1, sticky="ns", padx=6)
            self.wordLb.grid(row=0, column=2, sticky="w")
            self.sep3.grid(row=0, column=3, sticky="ns", padx=6)
            self.loadBt.grid(row=0, column=4, sticky="we")
        
        def _update_widget(self):
            """ 更新控件 """
            update = False
            if self.first:
                self.first = False
                update = True
            elif self.last != self.key:
                update = True
            # 执行刷新
            if update:
                # 高亮调整
                self.txtBrd.config(state="normal")
                self.txtBrd.tag_config(f"sen_{self.last}", background="white")
                self.txtBrd.tag_config(f"sen_{self.key}", background="yellow")
                self.txtBrd.config(state="disabled")
                # 信息栏调整
                if self.text_update:
                    self.wordLb.config(text=f"总字数: {len(self.text)}")
                    self.text_update = False
                self.placeLb.config(text=f"当前句子: {self.key + 1} / {len(self.stamp) + 1}")
                # 索引刷新
                self.last = self.key

        def _reload(self):
            """ 手动刷新 """
            # self.stop_load = True
            # while self.stop_load:  # 等待停止
            #     sleep(self.parent.monitor.configs["update"] // 2)
            # self.load()  # 万一有多个 after 怎么办? 自我删除重开? 算了...
            self.root.destroy()
            self.__init__(self.parent)  # 没除干净啊。算了，懒得 del。


    class GenerateBoard:
        """ PPT 生成准备界面 """

        class SlideHandler(tk.Frame):
            """
            单张参
            ---

            :param page: 页码
            :type page: int (from 0)

            :param get_config: 获取参
            :param reset_page: 修改页
            """

            def __init__(self, page: int, config: dict, parent, container, *args, **kwargs):
                """ 初始化
                :param page: 此页数 (非真实)
                :type page: int
                :param config: 初始参  **所有值，全都传**
                :type config: list[dict] -> 去看看 SpeechHandler.generate.slides (有不同)
                :param parent: 上一级
                :type parent: GenerateBoard
                :param container: 父容器
                :type container: tk.Frame """
                tk.Frame.__init__(self, container, *args, **kwargs)
                self.parent = parent

                self.pgeVr = tk.StringVar(self)
                self.reset_page(page)
                self.configs = config  # 存参数
                # 布容器
                self.config(highlightbackground=self.parent.monitor.colors["gray3"], highlightthickness=1)
                self._setup_widget()
                self._layout_widget()
                # 绑菜单
                self.bind("<Button-3>", self._post_menu)
                for widget in self.winfo_children():
                    widget.bind("<Button-3>", self._post_menu)
                # 列配置
                self.grid_columnconfigure([2, 5, 8], weight=1)
            
            def _setup_widget(self):
                """ 造控件 """
                self.boxwidth = 6
                # -> 弹菜单
                self.menu = tk.Menu(self, tearoff=False)
                self.menu.add_command(label="插入(I)", underline=3, command=lambda: self.parent.insert_slide(self.page, self.parent.default))
                self.menu.add_command(label="删除(D)", underline=3, command=lambda: self.parent.del_slide(self.page))
                self.menu.add_separator()
                self.menu.add_command(label="清除下方全部(C)", underline=7, command=lambda: self.parent.del_from(self.page + 1))

                # -> 通用态
                self.sep1 = ttk.Separator(self, orient="vertical")
                self.sep2 = ttk.Separator(self, orient="vertical")
                self.sep3 = ttk.Separator(self, orient="vertical")
                # 此页数
                self.pgeLb = tk.Label(self, textvariable=self.pgeVr, font=self.parent.monitor.fonts["text1_hold"])
                # 选类型
                self.typeCb = ttk.Combobox(self, state="readonly", width=self.boxwidth)
                self.typeCb["values"] = ("单句", "多句")
                self.typeCb.current(0 if self.configs["type"] == "single" else 1)
                self.typeCb.bind("<<ComboboxSelected>>", self._type_change)
                # 设标题
                self.ttVr = tk.StringVar(self, value=self.configs["title"] if self.configs["title"] else "默认")
                self.ttLb = tk.Label(self, text=f"标题:")
                self.ttEt = ttk.Entry(self, textvariable=self.ttVr, width=int(self.boxwidth * 2.5))
                # 显原句
                self.ansVr = tk.BooleanVar(self, value=self.configs["text"])
                self.ansRb = ttk.Checkbutton(self, text="追加呈现原句页", variable=self.ansVr, command=self._ans_change)

                # -> 单句态
                # 选单句
                self.keyLb = tk.Label(self, text="句子:")
                self.keySp = ttk.Spinbox(self, from_=1, to=len(self.parent.monitor.audio.info), increment=1, command=self._key_change, width=self.boxwidth)
                self.keySp.set(self.configs["key"] + 1)
                # 单遍开
                self.sngVr = tk.BooleanVar(self, value=self.configs["single"])
                self.sngRb = ttk.Checkbutton(self, text="读一遍音频", variable=self.sngVr, command=self._sng_change)
                # 多遍开
                self.rptVr = tk.BooleanVar(self, value=self.configs["repeat"]["on"])
                self.rptRb = ttk.Checkbutton(self, text="重复读音频", variable=self.rptVr, command=self._rpt_change)
                # 重复数
                self.rnmLb = tk.Label(self, text="重复次数:")
                self.rnmSp = ttk.Spinbox(self, from_=1, to=15, increment=1, command=self._rnm_change, width=self.boxwidth)
                self.rnmSp.set(self.configs["repeat"]["time"])
                # 遍间隔
                self.rpsLb = tk.Label(self, text="重复间隔:")
                self.rpsSp = ttk.Spinbox(self, from_=0.1, to=20.0, increment=0.1, command=self._rps_change, width=self.boxwidth)
                self.rpsSp.set(self.configs["repeat"]["silence"])
                # 单遍倍
                self.sspVr = tk.DoubleVar(self)
                self.sspLb = tk.Label(self, text="单遍倍速:")
                self.sspSp = ttk.Spinbox(self, from_=0.3, to=10.0, increment=0.1, command=self._ssp_change, textvariable=self.sspVr, width=self.boxwidth)
                self.sspSp.set(self.configs["speed"]["single"])
                # 多遍倍
                self.rspVr = tk.DoubleVar(self)
                self.rspLb = tk.Label(self, text="重复倍速:")
                self.rspSp = ttk.Spinbox(self, from_=0.3, to=10.0, increment=0.1, command=self._rsp_change, textvariable=self.rspVr, width=self.boxwidth)
                self.rspSp.set(self.configs["speed"]["repeat"])
                # 绑倍速
                self.bsrVr = tk.BooleanVar(self, value=False)
                self.bsrRb = ttk.Checkbutton(self, text="同步倍速调整", variable=self.bsrVr, command=self._bsr_change)
                
                # -> 多句态
                # 起点句
                self.frsLb = tk.Label(self, text="起点句:")
                self.frsSp = ttk.Spinbox(self, from_=1, to=len(self.parent.monitor.audio.info), increment=1, command=self._frs_change, width=self.boxwidth)
                self.frsSp.set(self.configs["side"][0] + 1)
                # 终点句
                self.edsLb = tk.Label(self, text="终点句:")
                self.edsSp = ttk.Spinbox(self, from_=1, to=len(self.parent.monitor.audio.info), increment=1, command=self._eds_change, width=self.boxwidth)
                if self.configs["side"][1] == None:
                    self.configs["side"][1] = len(self.parent.monitor.audio.info) - 1
                self.edsSp.set(self.configs["side"][1] + 1)
                self.esbVr = tk.BooleanVar(self, value=False)
                self.esbRb = ttk.Checkbutton(self, text="直到最后", variable=self.esbVr, command=self._esb_change)
                # 多句倍
                self.lspLb = tk.Label(self, text="倍速:")
                self.lspSp = ttk.Spinbox(self, from_=0.3, to=10.0, increment=0.1, command=self._lsp_change, width=self.boxwidth)
                self.lspSp.set(self.configs["speed"]["long"])
            
            def _layout_widget(self):
                """ 排控件 """
                # 单句排
                if self.configs["type"] == "single":
                    # 最左侧
                    self.pgeLb. grid(column=0,  row=0, columnspan=2, sticky="w" )
                    self.typeCb.grid(column=0,  row=1, columnspan=2, sticky="w" )
                    self.ttLb.  grid(column=0,  row=2, sticky="w"               )
                    self.ttEt.  grid(column=1,  row=2, sticky="w"               )
                    self.sep1.  grid(column=2,  row=0, rowspan=3,    sticky="ns", padx=(6, 12))
                    # 中间左
                    self.keyLb. grid(column=3,  row=0, sticky="w"               )
                    self.keySp. grid(column=4,  row=0, sticky="w"               )
                    self.sngRb. grid(column=3,  row=1, columnspan=2, sticky="w" )
                    self.rptRb. grid(column=3,  row=2, columnspan=2, sticky="w" )
                    self.sep2.  grid(column=5,  row=0, rowspan=3,    sticky="ns", padx=(6, 6) )
                    # 中间右
                    self.ansRb. grid(column=6,  row=0, columnspan=2, sticky="w" )
                    self.bsrRb. grid(column=6,  row=1, columnspan=2, sticky="w" )
                    self.sspLb. grid(column=6,  row=2, sticky="w"               )
                    self.sspSp. grid(column=7,  row=2, sticky="w"               )
                    self.sep3.  grid(column=8,  row=0, rowspan=3,    sticky="ns", padx=(6, 6) )
                    # 最右侧
                    self.rnmLb. grid(column=9,  row=0, sticky="w")
                    self.rnmSp. grid(column=10, row=0, sticky="w")
                    self.rpsLb. grid(column=9,  row=1, sticky="w")
                    self.rpsSp. grid(column=10, row=1, sticky="w")
                    self.rspLb. grid(column=9,  row=2, sticky="w")
                    self.rspSp. grid(column=10, row=2, sticky="w")
                # 多句排
                else:
                    # 最左侧
                    self.pgeLb. grid(column=0,  row=0, sticky="w"               )
                    self.typeCb.grid(column=0,  row=1, sticky="w"               )
                    self.sep1.  grid(column=2,  row=0, rowspan=2,    sticky="ns", padx=(6, 12))
                    # 中间左
                    self.ttLb.grid  (column=3,  row=0, sticky="w"               )
                    self.ttEt.grid  (column=4,  row=0, sticky="w"               )
                    self.keyLb.grid (column=3,  row=1, sticky="w"               )
                    self.keySp.grid (column=4,  row=1, sticky="w"               )
                    self.sep2.grid  (column=5,  row=0, rowspan=2,    sticky="ns", padx=(6, 6) )
                    # 中间右
                    self.ansRb.grid (column=6,  row=0, sticky="w"               )
                    self.esbRb.grid (column=6,  row=1, sticky="w"               )
                    self.sep3.grid  (column=8,  row=0, rowspan=2,    sticky="ns", padx=(6, 6) )
                    # 最右侧
                    self.frsLb.grid (column=9,  row=0, sticky="w")
                    self.frsSp.grid (column=10, row=0, sticky="w")
                    self.edsLb.grid (column=9,  row=1, sticky="w")
                    self.edsSp.grid (column=10, row=1, sticky="w")
            
            def _type_change(self, event):
                """ 类型变 """
                self.configs["type"] = "single" if self.typeCb.get() == "单句" else "long"
                # 全隐藏
                for widget in self.winfo_children():
                    widget.grid_forget()
                # 焕新生
                self._layout_widget()
            
            def _key_change(self):
                """ 单句变 """
                self.configs["key"] = int(self.keySp.get()) - 1

            def _ans_change(self):
                """ 显原句 """
                self.configs["text"] = self.ansVr.get()

            def _sng_change(self):
                """ 单遍改 """
                self.configs["single"] = value = self.sngVr.get()
                if value:
                    self.sspSp.config(state="normal")
                    self.bsrRb.config(value=0, state="normal")
                else:
                    self.sspSp.config(state="disabled")
                    self.bsrRb.config(value=0, state="disabled")
            
            def _rpt_change(self):
                """ 多遍改 """
                self.configs["repeat"]["on"] = value = self.rptVr.get()
                if value:
                    self.rspSp.config(state="normal")
                    self.bsrRb.config(value=0, state="normal")
                else:
                    self.rspSp.config(state="disabled")
                    self.bsrRb.config(value=0, state="disabled")
            
            def _rnm_change(self):
                """ 重复改 """
                self.configs["repeat"]["time"] = int(self.rnmSp.get())

            def _rps_change(self):
                """ 间隔改 """
                self.configs["repeat"]["silence"] = float(self.rpsSp.get())

            def _ssp_change(self):
                """ 单遍倍 """
                self.configs["speed"]["single"] = float(self.sspSp.get())

            def _rsp_change(self):
                """ 多遍倍 """
                self.configs["speed"]["repeat"] = float(self.rspSp.get())
            
            def _bsr_change(self):
                """ 绑倍速 """
                self.configs["speed"]["bind"] = value = self.bsrVr.get()
                if value:
                    self.rspSp.config(textvariable=self.sspVr)
                else:
                    self.rspSp.config(textvariable=self.rspVr)
            
            def _frs_change(self):
                """ 多句起 """
                self.configs["side"][0] = int(self.frsSp.get()) - 1

            def _eds_change(self):
                """ 多句终 """
                self.configs["side"][1] = int(self.edsSp.get()) - 1
            
            def _esb_change(self):
                """ 多句完 """
                self.configs["side"][2] = value = self.esbVr.get()
                if value:
                    self.edsSp.config(state="disabled")
                else:
                    self.edsSp.config(state="normal")
            
            def _lsp_change(self):
                """ 多句倍 """
                self.configs["speed"]["long"] = float(self.lspSp.get())
            
            def get_config(self) -> dict:
                """ 返回参
                :return: 格式参
                :rtype: dict -> 懒得写 """
                config = {
                    "title": None if self.configs["title"] == "默认" else self.configs["title"], 
                    "type": self.configs["type"], 
                    "text": self.configs["text"], 
                }
                # 单句参
                if self.configs["type"] == "single":
                    config.update({
                        "key": self.configs["key"], 
                        "single": self.configs["single"], 
                        "repeat": (int(self.configs["repeat"]["time"]) if self.configs["repeat"]["on"] else 0, self.configs["repeat"]["silence"]), 
                        "speed": None if (abs(self.configs["speed"]["single"] - 1.0) < 0.1 and abs(self.configs["speed"]["repeat"] - 1.0) < 0.1)\
                            else (self.configs["speed"]["single"] if self.config["speed"]["bind"] else (\
                                (None if (abs(self.configs["speed"]["single"] - 1.0) < 0.1) else self.configs["speed"]["single"]), \
                                    (None if (abs(self.configs["speed"]["repeat"] - 1.0) < 0.1) else self.configs["speed"]["repeat"]), )), 
                    })
                # 多句参
                else:
                    config.update({
                        "side": [self.configs["side"][0] if self.configs["side"][0] else None, None if self.configs["side"][2] else self.configs["side"][1]], 
                        "speed": None if abs(self.configs["speed"]["long"] - 1.0) < 0.1 else self.configs["speed"],  
                    })
                # 回参数
                return config.copy()  # 拷没用

            def reset_page(self, page: int):
                """ 修改页
                :param page: 新的页
                :type page: int """
                self.page = page
                self.pgeVr.set(f"第 {page + 1} 页")

            def _post_menu(self, event):
                """ 弹出菜单 """
                self.menu.post(event.x_root, event.y_root)

        default: dict = {  # 默认单页设置
            "type": "single", "title": "默认", "text": True, "key": 0, "single": True, "side": [0, None, False], 
            "repeat": {"on": True, "time": 3, "silence": 4.0, }, 
            "speed": {"single": 1.0, "repeat": 1.0, "bind": False, "long": 1.0}, 
        }

        def __init__(self, monitor):
            """ 初始化一个准备界面
            :param monitor: 最高父类
            :type monitor: Window """
            self.monitor = monitor
            # 窗口设置
            self.root = tk.Toplevel(self.monitor.root)
            self.root.minsize(600, 400)
            self.root.title("PPT 导出准备")
            # 控件设置
            self._setup_widget()
            self._layout_widget()
            self._first_load()
            # 运行窗口
            self.root.after(self.monitor.configs["update"], lambda: self._update_data())
            self.root.mainloop()

        def _setup_widget(self):
            """ 生成控件 """
            self.sep2 = ttk.Separator(self.root)
            self.sep3 = ttk.Separator(self.root)
            self.toolBar = tk.Frame(self.root)  # 顶部工具栏
            self.slideBrd = self.monitor.ScrollableFrame(self.root)  # PPT 呈现
            self.slideFrm = self.slideBrd.container  # 映射一下容器
            self.dataFrm = tk.Frame(self.root)  # 底部信息显示栏

            # 工具栏
            self.addBt = ttk.Button(self.toolBar, text="添加", width=5, command=lambda: self.add_slide(self.default))
            self.insBt = ttk.Button(self.toolBar, text="插入", width=5, command=self._ask_for_insert)
            self.sep1 = ttk.Separator(self.toolBar, orient="vertical")
            self.delBt = ttk.Button(self.toolBar, text="清空", width=5, command=lambda: self.del_from(0))
            self.genBt = ttk.Button(self.toolBar, text="生成", width=5, command=self.ask_export)

            # 呈现
            self.slides = []  # 盛放 SlideHandler

            # 信息栏
            self.pgeLb = tk.Label(self.dataFrm, text="页数: \\")
        
        def _layout_widget(self):
            """ 布局控件 """
            self.toolBar.pack(side="top", anchor="nw", fill="x")
            self.sep2.pack(side="top", fill="x")
            self.slideBrd.pack(side="top", anchor="nw", fill="both", expand=True)
            self.sep3.pack(side="top", fill="x")
            self.dataFrm.pack(side="bottom", anchor="nw", fill="x")
            # 工具栏
            self.addBt.pack(side="left", padx=4)
            self.insBt.pack(side="left")
            self.sep1.pack(side="left", padx=8, fill="y")
            self.delBt.pack(side="left")
            self.genBt.pack(side="right", padx=4)
            # 信息栏
            self.pgeLb.grid(row=0, column=0, sticky="w")
        
        def _update_data(self):
            """ 刷新底部信息栏 """
            self.pgeLb.config(text=f"页数: {len(self.slides)}")
            self.root.after(self.monitor.configs["update"], lambda: self._update_data())
        
        def _first_load(self):
            """ 初次加载 """
            config = deepcopy(self.default)
            # 单句添加
            for cnt in range(len(self.monitor.audio.info)):
                # 修正设置
                config["key"] = cnt
                self.add_slide(config)
            # 全文添加
            config["type"] = "long"
            self.add_slide(config)

        def del_slide(self, index: int):
            """ 删除索引
            :param index: 删除的索引
            :type index: int """
            # 删除
            self.slides[index].destroy()
            del self.slides[index]
            # 修正页面
            self._batch_page_edit(index - 1, 1)
        
        def del_from(self, index: int):
            """ 删除右侧页面
            :param index: 起始序列号
            :type index: int """
            for cnt in range(index, len(self.slides)):
                self.slides[cnt].destroy()
                del self.slides[cnt]

        def insert_slide(self, index: int, config: dict):
            """ 插入一页
            :param index: 插入位索引
            :type index: int
            :param config: 初始设置
            :type config: dict -> 懒得写 """
            # 插入
            self.slides.insert(index, self.SlideHandler(index, deepcopy(config), self, self.slideFrm))
            self.slides[index].pack(side="top", padx=2, pady=3, anchor="nw", fill="x")
            # 修改页面
            self._batch_page_edit(index + 1, 1)

        def add_slide(self, config: dict):
            """ 在末尾添加一页
            :param config: 初始设置
            :type config: dict -> 懒得写 """
            self.insert_slide(len(self.slides), config)

        def _batch_page_edit(self, from_: int, add_: int):
            """ 批量修改页码
            :param from_: 起始位点
            :type from_: int
            :param add_: 追加值
            :type add_: int """
            for slide in self.slides[from_:]:
                slide.reset_page(slide.page + add_)

        def _ask_for_insert(self):
            """ 询问并插入 """
            index = simpledialog.askinteger(title="询问", prompt="在第 X 页后插入，请输入 X 值:", minvalue=1, maxvalue=len(self.slides))
            if type(index) == int and index >= 1:
                index -= 1
                amount = simpledialog.askinteger(title="询问", prompt="插入 X 页，请输入 X 值:", minvalue=1, maxvalue=len(self.slides))
                for i in range(amount):
                    self.insert_slide(index := index + 1, self.default)

        def ask_export(self):
            """ 询问导出位置 """
            # 询问位置
            path = filedialog.asksaveasfilename(title="导出 PPT", defaultextension=".pptx", initialfile=f"听写-{self.monitor.audio.name}.pptx", filetypes=(("PPT 幻灯片", "*.pptx"), ))
            if path:
                self.export(path)
                self.root.destroy()
        
        def export(self, path: str):
            """ 生成并导出
            :param path: PPT 文件位置
            :type path: str """
            # 获取参数
            config = []
            for slide in self.slides:
                config += slide.get_config(), 
            # 执行导出
            self.monitor.audio.generate(config, path)


    class ConfigBoard:
        """
        设置修改窗口
        ---

        :param reset: 加载设置
        :param reload: 加载本地默认设置
        :param run: 应用设置
        :param save: 保存设置
        """

        def __init__(self, monitor):
            """ 初始化设置窗口
            :param monitor: 父控件
            :type monitor: Window """
            self.monitor = monitor
            # 窗口设置
            self.root = tk.Toplevel(self.monitor.root)
            self.root.title("设置")
            self.root.resizable(False, False)
            # print(self.root.winfo_height(), self.root.winfo_width())
            # 控件设置
            self._setup_widget()
            self._layout_widget()
            self.reset()  # 获取设置并显示
            # 运行窗口
            self.root.mainloop()
        
        def _setup_widget(self):
            """ 设置控件 """
            self.cfgFrm = tk.Frame(self.root)
            self.sep1 = ttk.Separator(self.root, orient="horizontal")
            self.hdlFrm = tk.Frame(self.root)

            # 设置修改区域
            self.sep3 = ttk.Separator(self.cfgFrm, orient="vertical")
            self.modVr = tk.StringVar(self.cfgFrm)  # 模型名称
            self.modLb = tk.Label(self.cfgFrm, text="模型:")
            self.modEt = ttk.Entry(self.cfgFrm, textvariable=self.modVr)
            self.fbdLb = tk.Label(self.cfgFrm, text="快进/快退时间(s):")  # 快进/快退
            self.fwdSb = ttk.Spinbox(self.cfgFrm, width=6, from_=0.1, to=60.0, increment=0.1, command=self._change_fwd)
            self.bwdSb = ttk.Spinbox(self.cfgFrm, width=6, from_=0.1, to=60.0, increment=0.1, command=self._change_bwd)
            self.fpsLb = tk.Label(self.cfgFrm, text="窗口刷新间隔(非帧率)(ms):")  # 刷新间隔
            self.fpsSb = ttk.Spinbox(self.cfgFrm, width=10, from_=10, to=10000, increment=10, command=self._change_fps)

            # 应用界面
            self.ldcBt = ttk.Button(self.hdlFrm, text="抹除修改", width=9, command=self.reset)
            self.rstBt = ttk.Button(self.hdlFrm, text="恢复默认", width=9, command=self.reload)
            self.sep2 = ttk.Separator(self.hdlFrm, orient="vertical")
            self.dftBt = ttk.Button(self.hdlFrm, text="设为默认", width=9, command=self.save)
            self.runBt = ttk.Button(self.hdlFrm, text="应用", width=5, command=self.run)
            self.kclBt = ttk.Button(self.hdlFrm, text="取消", width=5, command=self.root.destroy)

        def _layout_widget(self):
            """ 排布控件 """
            self.cfgFrm.pack(side="top", fill="both", anchor="nw", padx=(6, 2))
            self.sep1.pack(side="top", fill="x")
            self.hdlFrm.pack(side="bottom", fill="x", anchor="se")
            # 设置界面
            self.sep3. grid(row=0, column=1, sticky="ns", rowspan=3, padx=6)
            self.modLb.grid(row=0, column=0, sticky="e")
            self.modEt.grid(row=0, column=2, sticky="w",  columnspan=2)
            self.fbdLb.grid(row=1, column=0, sticky="e")
            self.fwdSb.grid(row=1, column=2, sticky="w")
            self.bwdSb.grid(row=1, column=3, sticky="w")
            self.fpsLb.grid(row=2, column=0, sticky="e")
            self.fpsSb.grid(row=2, column=2, sticky="w",  columnspan=2)
            # 应用界面
            self.ldcBt.grid(row=0, column=0, sticky="e",  padx=(0, 4))
            self.rstBt.grid(row=0, column=1, sticky="e")
            self.sep2.grid (row=0, column=2, sticky="ns", padx=12)
            self.dftBt.grid(row=0, column=3, sticky="e",  padx=(0, 6))
            self.runBt.grid(row=0, column=4, sticky="e",  padx=(0, 4))
            self.kclBt.grid(row=0, column=5, sticky="e",  padx=(0, 4))
        
        def _update_config(self):
            """ 加载控件的设置内容渲染 """
            self.modVr.set(self.config["model"])
            self.fwdSb.set(self.config["forward"])
            self.bwdSb.set(self.config["backward"])
            self.fpsSb.set(self.config["update"])
        
        def _change_fwd(self):
            """ 快进修改 """
            self.config["forward"] = float(self.fwdSb.get())

        def _change_bwd(self):
            """ 快退修改 """
            self.config["backward"] = float(self.bwdSb.get())

        def _change_fps(self):
            """ 刷新间隔修改 """
            self.config["update"] = int(self.fpsSb.get())
        
        def run(self):
            """ 应用设置 """
            self.monitor.config(deepcopy(self.config))
            self.root.destroy()
        
        def reset(self):
            """ 抹除设置修改 """
            self.config = deepcopy(self.monitor.configs)
            self._update_config()
        
        def reload(self):
            """ 加载本地默认设置 """
            with open(self.monitor.cfg_pt, mode="r", encoding="utf-8") as cfg_file:
                workpath = self.config["workpath"]
                self.config = load(cfg_file)
                self.config["workpath"] = workpath
            self._update_config()
        
        def save(self):
            """ 保存并同步设置 """
            config = deepcopy(self.config)
            if "workpath" in config.keys():
                del config["workpath"]
            with open(self.monitor.cfg_pt, mode="w", encoding="utf-8") as cfg_file:
                dump(config, cfg_file)
            self.run()


    ### 内部访问方法 ###

    def config(self, config: dict):
        """ 全局刷新设置
        :param config: 设置
        :type config: dict """
        self.configs = config
        # 同步到该去的地方
        self.audBrd.control.configs = self.configs  # 这也许没意义？但保险一点
        self.audio.load_config(self.configs)

    def config_load(self):
        """ 基础设置加载 """
        self.workpath = os.path.dirname(os.path.abspath("__file__"))  # 绝对路径
        # 配置项
        self.cfg_pt = os.path.join(self.workpath, "config.json")  # 配置文件路径
        with open(self.cfg_pt, mode="r", encoding="utf-8") as config_file:
            self.configs = load(config_file)
        self.configs["workpath"] = self.workpath
            
        # 颜色表
        self.colors = {
            "gray1": "#363636",       # 较深灰
            "gray2": "#4F4F4F",       # 中等灰
            "gray3": "#696969",       # 浅灰
            "gray11": "#1C1C1C",      # 深灰
            "firebrick1": "#FF3030",  # 耐火砖(红)
            "gold": "#FFD700",        # 金色
        }

        # 字体
        self.fonts = {
            "text1_hold": font.Font(family="黑体", size=12, weight="bold"),  # 加粗中等大字体
            # "text2_para": font.Font(family="Arial", size=10),  # 长段英文较小字体
            "text2_para": font.Font(family="Times New Roman", size=14),  # 长段英文较小字体
            # "text3_title": font.Font(family="等线", size=14),  # 正常大标题
            "text3_title": font.Font(family="Arial", size=14),  # 正常大标题
            "text4_context": font.Font(family="Times New Roman", size=24),   # 长文段较大字号
            "text5_text_ch": font.Font(family="宋体", size=18),   # 长文段较大字号
        }
    
    def _global_activate(self):
        """ 初次加载音频激活控件 """
        self.menu.entryconfig("导出(E)", state="normal")  # 激活菜单功能
        self.topBar.activate()  # 激活顶端按钮
        self.audBrd.activate()  # 激活播放按钮 & 调整按钮

    def global_update(self):
        """ 载入文件时全局刷新 """
        self.topBar.set_info()
        self.senBrd.load()
        self.audBrd.load()
        # 初次刷新？
        if self.first_load:
            self.first_load = False
            self._global_activate()


    ### 用户交互方法定义 ###

    def new_window(self, build_func=None):
        """ 新建窗口(可能会因为用户构建子类而出问题吧)
        :param build_func: 窗口创建函数 """
        if build_func == None:
            global set_window
            build_func = set_window
        # 拉入新的进程(不依赖)
        newWindowProcess = Process(target=build_func, daemon=False)
        newWindowProcess.start()

    def open_audio(self, path: str=None):
        """ 打开语料文件(音频)
        :param path: 已有路径，否则询问
        :type path: str = None """
        # 加载文件(防止为空)
        if path:
            self.audio.load(path)
            # 全局刷新
            self.global_update()
            return
        # 询问文件
        path = filedialog.askopenfilename(title="选择文件", filetypes=(("音频文件", "*.mp3 *.wav *.m4a *.flac"), ))
        # 防卡死
        self.root.after(1, lambda: self.open_audio(path))

    def open_video(self, path: str=None):
        """ 打开语料文件(视频)
        :param path: 已有路径，否则询问
        :type path: str = None """
        # 加载文件
        if path:
            self.audio.reset()
            self.audio.video_load(path)
            # 全局刷新
            self.global_update()
            return
        # 询问
        path = filedialog.askopenfilename(title="选择文件", filetypes=(("视频文件", ".mp4 *.mkv *.avi"), ))
        # 防卡死
        self.root.after(1, lambda: self.open_video(path))

    def export_slice(self):
        """ 导出音频切片 """
        sentence = simpledialog.askstring(title="询问导出范围", initialvalue="<|3|>, ALL", prompt="""请按下列要求输入导出的句子范围:
                                          1. 若需要全部句子, 仅输入 ALL
                                          2. 否则，请按照下列规则书写一段表达式:
                                          \t1. 使用 `,` 分隔每段信息
                                          \t2. 使用句子序号表示单个句子，如 `1` 表示第一句
                                          \t3. 若要表示连续的几句，请输入 `起始句->终点句(包含)`，如 `2->4` 表示第二到四句
                                          \t4. 若要在句间追加一段停顿，请输入 `|停顿时间|`，例如 `|1.3|` 表示停顿 1.3 秒
                                          \t5. 若要设置单句的倍速，请在该句开头紧接上 `x倍速|`，例如 `x1.5|3` 表示 1.5 倍速的第三句
                                          3. 对于 1 和 2 的表示形式，若想要设置默认句间停顿，请在开头插入一段 `<|停顿时间|>`
                                          4. 若想设置全局默认倍速，请在开头插入一段 `<x倍速>`
                                          5. 请严格遵守上述书写规范，但你可以在任意位置打上任意数量的空格，这是允许的\n
                                          例如，`<|3|>, 1, 3->4, |6|, 3->4` 表示按照 `第一句 + 停顿3秒 + 第三和第四句 + 停顿6秒 + 第三和第四句` 的方式组装音频""")
        try:
            # 按需生成
            sentences = sentence.replace(" ", "").split(sep=",")
            sep_time: int = 0
            gb_spd: int = 0  # 1 倍速
            for sen in [0, 1, 2]:
                if sen >= len(sentences):
                    break
                # 间隔时间
                if sentences[sen][:2] + sentences[sen][-2:] == "<||>":
                    sep_time: float = float(sentences[sen][2 : -2])
                    df_silent = AudioSegment.silent(1000 * sep_time)  # 默认静音片段
                    del sentences[sen]
                    sen -= 1
                # 全局倍速
                elif sentences[sen][:2] + sentences[sen][-1:] == "<x>":
                    gb_spd: float = float(sentences[sen][2 : -1])
                    del sentences[sen]
                    sen -= 1
            audio = AudioSegment.empty()  # 最终音频
            # 全部获取信息转译
            try:
                all_ = sentences.index("ALL", stop=3)
            except ValueError:
                pass
            else:
                if sep_time:
                    sentences = [str(i) for i in range(1, 1 + len(self.audio.info))]  # 视为每句排一遍
                else:
                    sentences = [f"1->{len(self.audio.info) + 1}"]
            # 分句添加
            cnt = 0
            for sen in sentences:
                # 静音添加
                if sen[:1] + sen[-1:] == "||":  # 追加静音
                    audio += AudioSegment.silent(1000 * float(sen[1 : -1]))
                    continue
                if cnt and not sentences[cnt - 1][0] in ["<", "|"]:  # 追加默认静音
                    audio += df_silent
                # 倍数设置
                spd = gb_spd
                if sen[0] == "x":
                    spd, sen = sen.split("|")
                    spd: float = float(spd)
                # 音频添加
                if "->" in sen:  # 追加片段
                    audio += speed_audio(self.audio.get_slice(*[int(i) - 1 for i in sen.split("->")]), spd) if spd else self.audio.get_slice(*[int(i) - 1 for i in sen.split("->")])
                else:  # 追加单句
                    audio += speed_audio(self.audio.get_slice(int(sen) - 1), spd) if spd else self.audio.get_slice(int(sen) - 1)
                cnt += 1
        except Exception as error:
            messagebox.showerror(title="报错", message=f"发生错误 \"{type(Exception)}\"\n错误信息如下:\n{error}")
        else:
            path = filedialog.asksaveasfilename(title=f"保存音频切片", defaultextension=".wav", initialfile=f"切片-{sentence if len(sentence) < 6 else f'{sentence[:1]}..{sentence[-3:]}'}.wav", filetypes=(("音频文件", "*.mp3 *.wav *.m4a *.flac"), ))
            if path:
                audio.export(path, format=os.path.splitext(path)[1:])  # 导出音频

    def export_txt(self):
        """ 导出全文 """
        path = filedialog.asksaveasfilename(title="导出全文", defaultextension=".txt", initialfile=f"文本识别-{self.audio.name}.txt", filetypes=(("文本文档", "*.txt"), ))
        if path:
            text = self.audio.get_text()[0]
            with open(path, mode="w", encoding="utf-8") as txt_file:
                txt_file.write(text)

    def generate(self):
        """ 生成 PPT 并导出 """
        generator = self.GenerateBoard(self)

    def help_board(self):
        """ 阅读帮助文档 """
        root = tk.Toplevel(self.root)
        root.minsize(400, 300)
        root.title("帮助文档 - README")
        text = tk.Text(root, font=self.fonts["text5_text_ch"], width=64, height=20)
        scrll = ttk.Scrollbar(root, command=text.yview)
        text.config(yscrollcommand=scrll.set)
        scrll.pack(side="right", fill="y")
        text.pack(side="left", fill="both", expand=True)
        with open("README.md", mode="r", encoding="utf-8") as help_file:
            text.insert("1.0", help_file.read())
        text.config(state="disabled")
        root.mainloop()

    def config_board(self):
        """ 呼出设置窗口 """
        setboard = self.ConfigBoard(self)

    def save_work(self):
        """ 保存项目 """
        path = filedialog.asksaveasfilename(title="保存此项目(对齐信息&音/视频)", defaultextension=".zip", initialfile=f"识别项目-{self.audio.name}.zip", filetypes=(("压缩包", "*.zip"), ))
        if path:
            self.audio.save(path)
    
    def load_work(self):
        """ 打开项目 """
        path = filedialog.askopenfilename(title="打开保存的项目", filetypes=(("压缩包", "*.zip"), ("你改后缀了？(点我无视后缀)", "*")))
        if path:
            self.audio.init_from(path)
            self.global_update()


    ### 窗口构建 ###

    def __init__(self):
        self.root = tk.Tk()
        self.root.withdraw()  # 隐藏
        self.config_load()  # 加载数据 & 绝对路径
        # 预载管理器和模型
        self.audio = SpeechHandler(self.configs)  # 反正拷贝开销低

        # 窗口预加载
        self.root.title("英语朗读音频处理器 - 提取、对齐文本，听写课件生成")
        self.root.minsize(800, 650)
        self.root.resizable(False, True)
        # self.root.bind("<Configure>", lambda e: print(e.width, e.height))
        # 控件和初始化设置
        self.first_load = True
        self._set_widget()
        self._pack_widget()
        # 关闭窗口函数覆写
        self.root.protocol("WM_DELETE_WINDOW", self.close)
    
    def _set_widget(self):
        """ 设置控件 """
        self.menu = self.Menu(self, tearoff=False)
        self.topBar = self.TopBar(self)  # 顶部状态栏
        self.spr1 = ttk.Separator(self.root)
        self.senBrd = self.SentenceBoard(self)  # 句子显示板
        self.spr2 = ttk.Separator(self.root)
        self.audBrd = self.AudioBoard(self)  # 音频操作板
    
    def _pack_widget(self):
        """ 排布控件 """
        self.topBar.pack(side="top", anchor="nw", fill="x", padx=4)
        self.spr1.pack(side="top", fill="x")
        self.senBrd.pack(side="top", fill="both", expand=True)
        self.spr2.pack(side="top", fill="x")
        self.audBrd.pack(side="bottom", anchor="nw", fill="x")
    
    def mainloop(self, *args, **kwards):
        """ 初始化窗口 """
        # 直接套 -_-
        self.root.deiconify()  # 出现
        self.root.mainloop(*args, **kwards)
    
    def close(self):
        """ 关闭窗口 """
        try:
            plt.close()  # 清除 matplotlib 内存占用
        finally:
            try:
                if self.audio.tmp_audio:
                    self.audio.clean()  # 清理 SpeechHandler 产生的临时音频文件的存储占用
            finally:
                self.root.destroy()  # 你必须是安全的！尽管我们在一个 Python 而非 Rust 环境里工作


def set_window():
    """ 创建窗口 """
    window = Window()
    window.mainloop()

if __name__ == "__main__":
    set_window()
